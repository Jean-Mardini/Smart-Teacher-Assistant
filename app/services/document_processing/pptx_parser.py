from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from .cleaners import clean_text
from .tables import extract_pptx_tables
from typing import List, Dict, Any
import os


def extract_pptx_text(filepath: str, document_id: str = "doc"):
    """
    Extract text, tables, and speaker notes from a PPTX file.

    Returns:
        pages:  list of {"page": N, "text": "..."}
        full_text: complete cleaned text
        tables: list of table dicts
        images: list of image dicts with path (no binary)
    """

    prs = Presentation(filepath)

    pages = []
    full_text = ""
    images = []
    image_counter = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_lines = []

        # Resolve slide title for image captions
        slide_title = ""
        for shape in slide.shapes:
            try:
                if shape.is_placeholder and shape.placeholder_format.idx in (0, 1):
                    slide_title = shape.text_frame.text.strip()
                    break
            except Exception:
                pass

        for shape in slide.shapes:

            # Text frames (titles, text boxes, content)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue

                    # Treat title placeholder as a heading
                    try:
                        if shape.is_placeholder and shape.placeholder_format.idx in (0, 1):
                            # idx 0 = title, idx 1 = center title
                            slide_lines.append(para_text.upper())
                        else:
                            slide_lines.append(para_text)
                    except Exception:
                        slide_lines.append(para_text)

            # Images (shape_type 13 == PICTURE)
            if shape.shape_type == 13:
                image_counter += 1
                img_id = f"img_{image_counter}"
                img_dir = os.path.join("outputs", "images", document_id)
                os.makedirs(img_dir, exist_ok=True)
                img_path = os.path.join(img_dir, f"{img_id}.png")
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                caption = slide_title if slide_title else shape.name
                images.append({
                    "image_id": img_id,
                    "page": slide_number,
                    "caption": caption,
                    "path": img_path,
                })

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_lines.append(f"[Notes: {notes_text}]")

        cleaned = clean_text("\n".join(slide_lines))
        pages.append({"page": slide_number, "text": cleaned})
        full_text += cleaned + "\n"

    tables = extract_pptx_tables(prs)

    return pages, full_text, tables, images
