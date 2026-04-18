"""
Table extraction for PDF, DOCX, and PPTX documents.
No LLM calls — purely structural/deterministic extraction.
"""

import re
import pdfplumber
from docx import Document
from pptx import Presentation
from typing import Any, Dict, List

# Caption patterns — English and French academic/professional documents
_CAPTION_RE = re.compile(
    r"\b(table|tableau|tab\.?|figure|fig\.?)\s*[.:–\-]?\s*\d+",
    re.IGNORECASE,
)



# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _flatten_rows(rows: List[List[str]]) -> str:
    """Flatten a 2-D cell matrix to pipe-separated plain text for RAG.

    The first row is treated as a header and followed by a separator line.
    Entirely empty rows are dropped.
    """
    non_empty = [row for row in rows if any(c.strip() for c in row)]
    if not non_empty:
        return ""
    lines: List[str] = []
    for idx, row in enumerate(non_empty):
        # Collapse internal newlines / tabs within a cell to a single space
        # so multi-line cell content doesn't break the pipe-delimited format.
        line = " | ".join(" ".join(c.split()) for c in row)
        lines.append(line)
        if idx == 0 and len(non_empty) > 1:
            lines.append("-" * min(len(line), 80))
    return "\n".join(lines)


def _is_valid_table(rows: List[List[str]]) -> bool:
    """Return False for degenerate tables.

    A table must have at least 2 rows, at least 2 columns, and at least
    one non-empty cell — otherwise it is likely a mis-detected text block.
    """
    if len(rows) < 2:
        return False
    if max((len(r) for r in rows), default=0) < 2:
        return False
    return any(c.strip() for row in rows for c in row)


def _dedup_merged(row_cells) -> List[str]:
    """Return cell texts from a python-docx row with merged cells collapsed.

    python-docx exposes horizontally-merged cells as repeated references to the
    same underlying XML element (``_tc``).  Tracking ``id(cell._tc)`` skips them.
    """
    seen: set = set()
    result: List[str] = []
    for cell in row_cells:
        tc_id = id(cell._tc)
        if tc_id not in seen:
            seen.add(tc_id)
            result.append(cell.text.strip())
    return result


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _pdf_caption(page, bbox: tuple, gap: int = 60) -> str:
    """Return the heading or caption text found immediately above the table.

    Crops a strip of *gap* points above the table and returns the last
    non-empty line in that strip (i.e. the line closest to the table).
    This captures section headings like "2. The Themes (Choose One)" that
    appear directly above a table but don't match a "Table N" pattern.

    Falls back to checking a narrow strip below the table for an explicit
    "Table N" / "Figure N" style caption when nothing is found above.
    Returns an empty string when nothing is found in either location.
    """
    x0, y0, x1, y1 = bbox
    w, h = float(page.width), float(page.height)

    # Strip above the table — return the line closest to the table top edge
    if y0 > gap:
        above = page.crop((0.0, max(0.0, y0 - gap), w, y0))
        text = (above.extract_text() or "").strip()
        if text:
            lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
            if lines:
                return lines[-1]

    # Strip below the table — only for explicit "Table N" / "Figure N" captions
    if y1 < h - gap:
        below = page.crop((0.0, y1, w, min(h, y1 + gap)))
        text = (below.extract_text() or "").strip()
        if text and _CAPTION_RE.search(text):
            return text.split("\n")[0].strip()

    return ""


def extract_pdf_tables(filepath: str) -> List[Dict[str, Any]]:
    """Extract tables from a PDF using pdfplumber.

    Detection strategy (per page):
    1. Line-based detection (bordered tables) via ``page.find_tables()``.
    2. Text-alignment fallback for borderless/grid-less tables.

    Returns a list of dicts: ``{table_id, page, caption, text}``
    """
    results: List[Dict[str, Any]] = []
    counter = 0

    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            page_num = page.page_number

            # Line-based detection only (requires actual grid lines / borders).
            # The text-alignment fallback is intentionally absent — it mis-fires
            # on multi-column text layouts that have no real table structure.
            found = page.find_tables()

            for tbl in found:
                raw = tbl.extract()
                if not raw:
                    continue

                rows = [[cell or "" for cell in row] for row in raw if row]
                if not _is_valid_table(rows):
                    continue

                flat = _flatten_rows(rows)
                if not flat:
                    continue

                counter += 1
                caption = _pdf_caption(page, tbl.bbox) or f"Table {counter}"
                results.append({
                    "table_id": f"tbl_{counter}",
                    "page": page_num,
                    "caption": caption,
                    "text": flat,
                })

    return results


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_docx_tables(doc: Document) -> List[Dict[str, Any]]:
    """Extract tables from an already-opened python-docx ``Document``.

    Merged cells are collapsed so their text does not appear twice in a row.
    DOCX has no reliable page numbers; all tables are reported on page 1.

    Returns a list of dicts: ``{table_id, page, caption, text}``
    """
    results: List[Dict[str, Any]] = []

    for idx, table in enumerate(doc.tables, start=1):
        rows = [_dedup_merged(row.cells) for row in table.rows]
        if not _is_valid_table(rows):
            continue
        flat = _flatten_rows(rows)
        if not flat:
            continue
        results.append({
            "table_id": f"tbl_{idx}",
            "page": 1,
            "caption": f"Table {idx}",
            "text": flat,
        })

    return results


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_pptx_tables(prs: Presentation) -> List[Dict[str, Any]]:
    """Extract tables from an already-opened python-pptx ``Presentation``.

    Each table is linked to its slide number (used as page number).
    When the slide has a title placeholder, it is prepended to the caption.

    Returns a list of dicts: ``{table_id, page, caption, text}``
    """
    results: List[Dict[str, Any]] = []
    counter = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Collect slide title for richer captions
        slide_title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                if shape.is_placeholder and shape.placeholder_format.idx in (0, 1):
                    slide_title = shape.text_frame.text.strip()
                    break
            except Exception:
                pass

        for shape in slide.shapes:
            if not shape.has_table:
                continue

            rows = [
                [cell.text.strip() for cell in row.cells]
                for row in shape.table.rows
            ]
            if not _is_valid_table(rows):
                continue
            flat = _flatten_rows(rows)
            if not flat:
                continue

            counter += 1
            caption = (
                f"Slide {slide_num} — {slide_title}"
                if slide_title
                else f"Slide {slide_num} Table {counter}"
            )
            results.append({
                "table_id": f"tbl_{counter}",
                "page": slide_num,
                "caption": caption,
                "text": flat,
            })

    return results
