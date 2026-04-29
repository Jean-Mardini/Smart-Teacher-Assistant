"""Flexible Grader â€” rubric generation, reference-grounded scoring, and exports (Kristy).

This module is the **authoritative implementation** of Kristy's capstone evaluation work: Groq
(OpenAI-compatible client), PDF/DOCX/PPTX text extraction, lightweight reference retrieval,
batch grading, and TXT/DOCX/HTML reports. Configuration and history live under ``data/evaluation/``.

- **Streamlit UI:** ``streamlit_flexible_grader.py`` (imports this module).
- **HTTP API:** ``app/api/routers/evaluation.py`` calls into here.

Team-facing wrappers: :mod:`rubrics`, :mod:`grading`, :mod:`feedback`, :mod:`analytics` re-export
the relevant entry points for clearer imports elsewhere in the codebase.
"""

from __future__ import annotations

import csv
import hashlib
import html
import io
import json
import os
import re
import time
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import logging

from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from app.services.llm import groq_client
from app.services.knowledge.chunking import chunk_document, chunk_text as knowledge_chunk_text
from app.services.knowledge.indexing_pipeline import get_local_document_by_id
from app.services.knowledge.retrieval import Retriever
from app.services.knowledge.vector_store import LocalVectorStore
from app.storage.files import get_evaluation_dir

logger = logging.getLogger(__name__)

APP_TITLE = "Flexible Grader"
DEFAULT_MODEL = "llama-3.3-70b-versatile"
UPLOAD_TYPES = ["pdf", "docx", "pptx", "txt", "md", "json", "csv", "html", "rtf"]


def config_path() -> str:
    return str(get_evaluation_dir() / "config.json")


def presets_path() -> str:
    return str(get_evaluation_dir() / "rubric_presets.json")


def history_path() -> str:
    return str(get_evaluation_dir() / "history.jsonl")


def temp_uploads_dir() -> Path:
    path = get_evaluation_dir() / "tmp"
    path.mkdir(parents=True, exist_ok=True)
    return path


# =========================================================
# Storage
# =========================================================
def safe_read_json(path: str) -> Dict[str, Any]:
    if not os.path.exists(path):
        return {}
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception:
        logger.warning("Could not read JSON at %s", path, exc_info=True)
        return {}


def safe_write_json(path: str, data: Dict[str, Any]) -> None:
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        logger.error("Could not write JSON to %s", path, exc_info=True)


def append_history(record: Dict[str, Any]) -> None:
    try:
        with open(history_path(), "a", encoding="utf-8") as f:
            f.write(json.dumps(record, ensure_ascii=False) + "\n")
    except Exception:
        logger.error("Could not append to history", exc_info=True)


def load_history(limit: int = 100) -> List[Dict[str, Any]]:
    hp = history_path()
    if not os.path.exists(hp):
        return []
    rows = []
    try:
        with open(hp, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    rows.append(json.loads(line))
                except Exception:
                    continue
    except Exception:
        return []
    rows.reverse()
    return rows[:limit]


def parse_history_timestamp(value: str) -> Optional[datetime]:
    text = str(value or "").strip()
    if not text:
        return None
    for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d"):
        try:
            return datetime.strptime(text, fmt)
        except Exception:
            continue
    return None


def _record_percent(record: Dict[str, Any]) -> float:
    try:
        score = float(record.get("overall_score", 0) or 0)
        out_of = float(record.get("overall_out_of", 0) or 0)
    except Exception:
        return 0.0
    if out_of <= 0:
        return 0.0
    return round((score / out_of) * 100.0, 2)


def _normalize_history_type(value: str) -> str:
    raw = str(value or "").strip().lower()
    if raw in {"single", "single_submission"}:
        return "single"
    if raw in {"batch", "batch_submission"}:
        return "batch_submission"
    return "single"


def history_batch_id(batch_name: str, records: List[Dict[str, Any]]) -> str:
    seed = "|".join(
        [
            batch_name.strip(),
            now_iso(),
            *[str(record.get("title", "")).strip() for record in records],
        ]
    )
    return f"batch-{short_hash(seed, 12)}"


def filter_history_records(
    records: List[Dict[str, Any]],
    date_from: str = "",
    date_to: str = "",
    search: str = "",
    history_type: str = "all",
) -> List[Dict[str, Any]]:
    start = parse_history_timestamp(date_from)
    end = parse_history_timestamp(date_to)
    if end and len(str(date_to or "").strip()) <= 10:
        end = end.replace(hour=23, minute=59, second=59)

    search_tokens = [token for token in re.split(r"\s+", str(search or "").strip().lower()) if token]
    normalized_type = str(history_type or "all").strip().lower()

    out: List[Dict[str, Any]] = []
    for record in records:
        stamp = parse_history_timestamp(str(record.get("timestamp", "")))
        if start and (stamp is None or stamp < start):
            continue
        if end and (stamp is None or stamp > end):
            continue

        record_type = _normalize_history_type(str(record.get("history_type", "")))
        if normalized_type == "single" and record_type != "single":
            continue
        if normalized_type in {"batch", "batch_submission"} and record_type != "batch_submission":
            continue

        if search_tokens:
            haystack = " ".join(
                [
                    str(record.get("title", "")),
                    str(record.get("batch_name", "")),
                    str(record.get("id", "")),
                    str(record.get("timestamp", "")),
                ]
            ).lower()
            if not all(token in haystack for token in search_tokens):
                continue

        out.append(record)
    return out


def build_history_batches(records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    grouped: Dict[str, List[Dict[str, Any]]] = {}
    for record in records:
        batch_id = str(record.get("batch_id", "")).strip()
        if not batch_id:
            continue
        grouped.setdefault(batch_id, []).append(record)

    summaries: List[Dict[str, Any]] = []
    for batch_id, batch_records in grouped.items():
        sorted_records = sorted(
            batch_records,
            key=lambda row: float(row.get("overall_score", 0) or 0),
            reverse=True,
        )
        percents = [_record_percent(record) for record in sorted_records]
        created_at = sorted_records[0].get("batch_created_at") or sorted_records[0].get("timestamp", "")
        batch_name = sorted_records[0].get("batch_name") or "Batch"
        total_scores = [float(record.get("overall_score", 0) or 0) for record in sorted_records]
        total_out_of = [float(record.get("overall_out_of", 0) or 0) for record in sorted_records]

        summaries.append(
            {
                "batch_id": batch_id,
                "batch_name": batch_name,
                "created_at": created_at,
                "submission_count": len(sorted_records),
                "average_percent": round(sum(percents) / len(percents), 2) if percents else 0.0,
                "highest_percent": max(percents) if percents else 0.0,
                "lowest_percent": min(percents) if percents else 0.0,
                "average_score": round(sum(total_scores) / len(total_scores), 2) if total_scores else 0.0,
                "average_out_of": round(sum(total_out_of) / len(total_out_of), 2) if total_out_of else 0.0,
                "records": sorted_records,
            }
        )

    summaries.sort(key=lambda row: str(row.get("created_at", "")), reverse=True)
    return summaries


def build_history_stats(records: List[Dict[str, Any]], batches: List[Dict[str, Any]]) -> Dict[str, Any]:
    percents = [_record_percent(record) for record in records]
    return {
        "total_records": len(records),
        "single_records": sum(1 for record in records if _normalize_history_type(record.get("history_type", "")) == "single"),
        "batch_records": sum(1 for record in records if _normalize_history_type(record.get("history_type", "")) == "batch_submission"),
        "batch_count": len(batches),
        "average_percent": round(sum(percents) / len(percents), 2) if percents else 0.0,
        "highest_percent": max(percents) if percents else 0.0,
        "lowest_percent": min(percents) if percents else 0.0,
    }


def load_history_view(
    limit: int = 100,
    date_from: str = "",
    date_to: str = "",
    search: str = "",
    history_type: str = "all",
) -> Dict[str, Any]:
    records = load_history(limit=limit)
    filtered = filter_history_records(records, date_from=date_from, date_to=date_to, search=search, history_type=history_type)
    batches = build_history_batches(filtered)
    stats = build_history_stats(filtered, batches)
    return {"records": filtered, "batches": batches, "stats": stats}


def save_history_records(records: List[Dict[str, Any]]) -> None:
    try:
        with open(history_path(), "w", encoding="utf-8") as f:
            for record in records:
                f.write(json.dumps(record, ensure_ascii=False) + "\n")
    except Exception:
        logger.error("Could not save history records", exc_info=True)


def _coerce_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except Exception:
        return default


def _coerce_int_points(value: Any) -> int:
    """Rubric max points as a non-negative int (handles null / missing / strings from JSON)."""
    if value is None:
        return 0
    try:
        f = float(value)
        if f != f or not f < float("inf"):  # NaN or not finite
            return 0
        return max(0, int(f))
    except (TypeError, ValueError, OverflowError):
        return 0


def normalize_result_record(record: Dict[str, Any]) -> Dict[str, Any]:
    normalized = dict(record or {})
    raw_items = normalized.get("items_results", []) or []
    items_results: List[Dict[str, Any]] = []
    total_score = 0.0
    total_out_of = 0.0

    for raw_item in raw_items:
        item = dict(raw_item or {})
        max_points = max(0.0, _coerce_float(item.get("points", 0), 0.0))
        earned_points = _coerce_float(item.get("earned_points", 0), 0.0)
        earned_points = max(0.0, min(max_points, earned_points))

        suggestions = item.get("suggestions", [])
        if isinstance(suggestions, str):
            suggestions = [line.strip() for line in suggestions.splitlines() if line.strip()]
        elif not isinstance(suggestions, list):
            suggestions = []

        evidence = item.get("evidence", [])
        if not isinstance(evidence, list):
            evidence = []

        item["points"] = round(max_points, 2)
        item["earned_points"] = round(earned_points, 2)
        item["rationale"] = str(item.get("rationale", "") or "").strip()
        item["suggestions"] = [str(line or "").strip() for line in suggestions if str(line or "").strip()]
        item["evidence"] = [entry for entry in evidence if isinstance(entry, dict)]
        items_results.append(item)

        total_score += earned_points
        total_out_of += max_points

    normalized["items_results"] = items_results
    normalized["overall_score"] = round(total_score, 2)
    normalized["overall_out_of"] = int(round(total_out_of))
    normalized["score_percent"] = _record_percent(normalized)
    normalized["manual_reviewed"] = bool(normalized.get("manual_reviewed", False))
    normalized["updated_at"] = now_iso()
    return normalized


def update_history_record(record_id: str, updated_record: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    target_id = str(record_id or "").strip()
    if not target_id:
        return None

    hp = history_path()
    if not os.path.exists(hp):
        return None

    rows: List[Dict[str, Any]] = []
    try:
        with open(hp, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    rows.append(json.loads(line))
                except Exception:
                    continue
    except Exception:
        logger.error("Could not read history for update", exc_info=True)
        return None

    saved: Optional[Dict[str, Any]] = None
    for index, row in enumerate(rows):
        if str(row.get("id", "")).strip() != target_id:
            continue
        merged = dict(row)
        merged.update(dict(updated_record or {}))
        merged["id"] = target_id
        merged.setdefault("timestamp", row.get("timestamp", now_iso()))
        merged["manual_reviewed"] = True
        saved = normalize_result_record(merged)
        rows[index] = saved
        break

    if saved is None:
        return None

    save_history_records(rows)
    return saved


def clear_history() -> None:
    try:
        hp = history_path()
        if os.path.exists(hp):
            os.remove(hp)
    except Exception:
        logger.error("Could not clear history", exc_info=True)


def load_config() -> Dict[str, Any]:
    return safe_read_json(config_path())


def save_config(cfg: Dict[str, Any]) -> None:
    safe_write_json(config_path(), cfg)


def load_presets() -> Dict[str, Any]:
    return safe_read_json(presets_path())


def save_presets(presets: Dict[str, Any]) -> None:
    safe_write_json(presets_path(), presets)


def now_iso() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def short_hash(text: str, n: int = 10) -> str:
    return hashlib.sha256((text or "").encode("utf-8", errors="ignore")).hexdigest()[:n]


def clean_spaces(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def escape_html(value: Any) -> str:
    return html.escape(str(value or ""), quote=True)


# =========================================================
# File extraction
# =========================================================
def extract_text(file_obj: Any, filename: str) -> str:
    try:
        if hasattr(file_obj, "seek"):
            file_obj.seek(0)
        data = file_obj.read() if hasattr(file_obj, "read") else b""
    except Exception:
        logger.warning("Could not read upload %s", filename, exc_info=True)
        data = b""

    if isinstance(data, str):
        return data.strip()
    return extract_text_from_bytes(data, filename)


def combine_uploaded_texts(files: List[Any], manual_text: str) -> str:
    parts: List[str] = []
    for f in files or []:
        try:
            fname = getattr(f, "name", None) or getattr(f, "filename", None) or "upload"
            txt = extract_text(f, fname)
            if txt.strip():
                parts.append(f"=== FILE: {fname} ===\n{txt}")
        except Exception:
            logger.warning("Could not combine uploaded text", exc_info=True)

    if (manual_text or "").strip():
        parts.append(f"=== MANUAL TEXT ===\n{manual_text.strip()}")

    return "\n\n".join(parts).strip()


def decode_text_bytes(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(enc, errors="ignore").strip()
        except Exception:
            continue
    return data.decode("utf-8", errors="ignore").strip()


def extract_csv_bytes(data: bytes) -> str:
    text = decode_text_bytes(data)
    if not text:
        return ""

    try:
        reader = csv.reader(io.StringIO(text))
        lines = ["[CSV TABLE]"]
        for i, row in enumerate(reader, start=1):
            cells = [clean_spaces(c) for c in row]
            if any(cells):
                lines.append(f"Row {i}: " + " | ".join(cells))
        return "\n".join(lines).strip()
    except Exception:
        return text


def extract_json_bytes(data: bytes) -> str:
    text = decode_text_bytes(data)
    if not text:
        return ""
    try:
        parsed = json.loads(text)
        return "[JSON]\n" + json.dumps(parsed, ensure_ascii=False, indent=2)
    except Exception:
        return text


def extract_html_bytes(data: bytes) -> str:
    text = decode_text_bytes(data)
    if not text:
        return ""
    text = re.sub(r"(?is)<script.*?>.*?</script>", " ", text)
    text = re.sub(r"(?is)<style.*?>.*?</style>", " ", text)
    text = re.sub(r"(?i)<br\s*/?>", "\n", text)
    text = re.sub(r"(?i)</(p|div|h1|h2|h3|li|tr|table)>", "\n", text)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    lines = [clean_spaces(line) for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def extract_rtf_bytes(data: bytes) -> str:
    text = decode_text_bytes(data)
    text = re.sub(r"\\par[d]?", "\n", text)
    text = re.sub(r"\\tab", "\t", text)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ")
    lines = [clean_spaces(line) for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def extract_with_processing_pipeline(data: bytes, filename: str) -> str:
    """Use the existing app processing pipeline for PDF/DOCX/PPTX."""
    try:
        from app.services.document_processing.pipeline import process_document
    except Exception:
        logger.warning("Processing pipeline import failed; using fallback extraction.", exc_info=True)
        return ""

    suffix = Path(filename or "upload").suffix.lower()
    if suffix not in {".pdf", ".docx", ".pptx"}:
        return ""

    temp_path = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(data)
            temp_path = tmp.name

        parsed = process_document(temp_path, document_number=1)
        parts: List[str] = []

        title = getattr(parsed, "title", "") or Path(filename).stem
        parts.append(f"[DOCUMENT TITLE]\n{title}")

        full_text = getattr(parsed, "full_text", "") or ""
        if full_text.strip():
            parts.append(f"[FULL TEXT]\n{full_text.strip()}")

        sections = getattr(parsed, "sections", []) or []
        for sec in sections:
            heading = sec.get("heading", "Section")
            level = sec.get("level", "")
            page_start = sec.get("page_start", "")
            text = sec.get("text", "")
            if text.strip():
                parts.append(f"[SECTION level={level} page={page_start}] {heading}\n{text.strip()}")

        tables = getattr(parsed, "tables", []) or []
        for tbl in tables:
            caption = tbl.get("caption", "Table")
            page = tbl.get("page", "")
            text = tbl.get("text", "")
            if text.strip():
                parts.append(f"[TABLE page={page}] {caption}\n{text.strip()}")

        images = getattr(parsed, "images", []) or []
        for img in images:
            caption = img.get("caption", "Image")
            page = img.get("page", "")
            path = img.get("path", "")
            parts.append(f"[IMAGE page={page}] {caption} {path}".strip())

        return "\n\n".join(parts).strip()
    except Exception:
        logger.warning("Processing pipeline extraction failed for %s", filename, exc_info=True)
        return ""
    finally:
        if temp_path:
            try:
                os.remove(temp_path)
            except Exception:
                pass


def extract_pdf_fallback(data: bytes) -> str:
    try:
        reader = PdfReader(io.BytesIO(data))
        parts: List[str] = []
        for page in reader.pages:
            text = (page.extract_text() or "").strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts).strip()
    except Exception:
        logger.warning("PDF fallback extraction failed", exc_info=True)
        return ""


def extract_docx_fallback(data: bytes) -> str:
    try:
        doc = Document(io.BytesIO(data))
        return "\n".join(
            paragraph.text.strip()
            for paragraph in doc.paragraphs
            if (paragraph.text or "").strip()
        ).strip()
    except Exception:
        logger.warning("DOCX fallback extraction failed", exc_info=True)
        return ""


def extract_pptx_fallback(data: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(data))
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if (text or "").strip():
                    parts.append(text.strip())
        return "\n".join(parts).strip()
    except Exception:
        logger.warning("PPTX fallback extraction failed", exc_info=True)
        return ""


def extract_text_from_bytes(data: bytes, filename: str) -> str:
    suffix = Path(filename or "").suffix.lower()

    if suffix == ".pdf":
        text = extract_with_processing_pipeline(data, filename)
        if text.strip():
            return text.strip()
        return extract_pdf_fallback(data)

    if suffix == ".docx":
        text = extract_with_processing_pipeline(data, filename)
        if text.strip():
            return text.strip()
        return extract_docx_fallback(data)

    if suffix == ".pptx":
        text = extract_with_processing_pipeline(data, filename)
        if text.strip():
            return text.strip()
        return extract_pptx_fallback(data)

    if suffix == ".csv":
        return extract_csv_bytes(data)
    if suffix == ".json":
        return extract_json_bytes(data)
    if suffix in {".html", ".htm"}:
        return extract_html_bytes(data)
    if suffix == ".rtf":
        return extract_rtf_bytes(data)

    return decode_text_bytes(data)


def extract_files_from_zip(zip_file: Any) -> List[Dict[str, str]]:
    out: List[Dict[str, str]] = []
    try:
        if hasattr(zip_file, "seek"):
            zip_file.seek(0)
        raw = zip_file.read() if hasattr(zip_file, "read") else zip_file.getvalue()
        with zipfile.ZipFile(io.BytesIO(raw), "r") as zf:
            for member in zf.infolist():
                if member.is_dir():
                    continue
                name = member.filename
                ext = Path(name).suffix.lower().lstrip(".")
                if ext not in UPLOAD_TYPES:
                    continue
                text = extract_text_from_bytes(zf.read(member), name)
                if text.strip():
                    out.append({"name": name, "text": text})
    except Exception:
        logger.warning("Could not extract ZIP", exc_info=True)
    return out


def normalize_uploaded_submissions(files: List[Any], zip_file: Optional[Any]) -> List[Dict[str, str]]:
    subs: List[Dict[str, str]] = []
    for f in files or []:
        try:
            name = getattr(f, "name", None) or getattr(f, "filename", None) or "upload"
            txt = extract_text(f, name)
            if txt.strip():
                subs.append({"name": name, "text": txt})
        except Exception:
            logger.warning("Could not parse uploaded submission", exc_info=True)

    if zip_file is not None:
        subs.extend(extract_files_from_zip(zip_file))

    return dedupe_named_texts(subs)


def parse_uploaded_files(files: List[Any]) -> List[Dict[str, str]]:
    parsed: List[Dict[str, str]] = []

    for upload in files or []:
        try:
            filename = getattr(upload, "filename", None) or getattr(upload, "name", None) or "upload"
            raw = upload.file.read() if hasattr(upload, "file") else upload.read()
        except Exception:
            logger.warning("Could not read upload", exc_info=True)
            continue

        if Path(filename).suffix.lower() == ".zip":
            try:
                with zipfile.ZipFile(io.BytesIO(raw), "r") as zf:
                    for member in zf.infolist():
                        if member.is_dir():
                            continue
                        name = member.filename
                        ext = Path(name).suffix.lower().lstrip(".")
                        if ext not in UPLOAD_TYPES:
                            continue
                        text = extract_text_from_bytes(zf.read(member), name)
                        if text.strip():
                            parsed.append({"name": name, "text": text})
            except Exception:
                logger.warning("Could not parse ZIP upload", exc_info=True)
            continue

        ext = Path(filename).suffix.lower().lstrip(".")
        if ext not in UPLOAD_TYPES:
            continue

        text = extract_text_from_bytes(raw, filename)
        if text.strip():
            parsed.append({"name": filename, "text": text})

    return dedupe_named_texts(parsed)


def dedupe_named_texts(items: List[Dict[str, str]]) -> List[Dict[str, str]]:
    out: List[Dict[str, str]] = []
    seen = set()
    for item in items:
        name = str(item.get("name", "upload"))
        text = str(item.get("text", ""))
        key = (name, short_hash(text, 16))
        if key in seen:
            continue
        seen.add(key)
        out.append({"name": name, "text": text})
    return out


def get_document_text(document_id: str) -> Dict[str, str]:
    doc = get_local_document_by_id(document_id)
    if doc is None:
        return {"document_id": document_id, "title": document_id, "text": ""}
    return {
        "document_id": doc.document_id,
        "title": doc.title,
        "text": (doc.full_text or "").strip(),
    }


def compose_text_from_sources(manual_text: str = "", document_ids: Optional[List[str]] = None) -> Dict[str, Any]:
    parts: List[str] = []
    documents: List[Dict[str, str]] = []

    for document_id in document_ids or []:
        resolved = get_document_text(document_id)
        text = resolved.get("text", "").strip()
        if not text:
            continue
        documents.append({"document_id": resolved["document_id"], "title": resolved["title"]})
        parts.append(f"=== DOCUMENT: {resolved['title']} ({resolved['document_id']}) ===\n{text}")

    if (manual_text or "").strip():
        parts.append(f"=== MANUAL TEXT ===\n{manual_text.strip()}")

    return {"text": "\n\n".join(parts).strip(), "documents": documents}


# =========================================================
# LLM  (delegates to the shared groq_client)
# =========================================================
def llm_json(
    payload: Dict[str, Any],
    model: Optional[str] = None,
    temperature: float = 0.15,
) -> Dict[str, Any]:
    return groq_client.call_llm_json_payload(payload, model=model, temperature=temperature)


# =========================================================
# Models
# =========================================================
def sanitize_assignment_item(raw: Dict[str, Any]) -> Dict[str, Any]:
    grounding = str(raw.get("grounding", "ai")).strip().lower()
    if grounding not in ["ai", "reference", "hybrid"]:
        grounding = "ai"

    pts = raw.get("points", 0)
    try:
        pts = int(pts)
    except Exception:
        pts = 0

    return {
        "item_origin": "assignment",
        "name": str(raw.get("name", "New Item")).strip() or "New Item",
        "description": str(raw.get("description", "")).strip(),
        "points": max(0, pts),
        "grounding": grounding,
        "expected_answer": "",
        "mode": "",
    }


def sanitize_teacher_key_item(raw: Dict[str, Any]) -> Dict[str, Any]:
    mode = str(raw.get("mode", "conceptual")).strip().lower()
    if mode not in ["exact", "conceptual"]:
        mode = "conceptual"

    grounding = str(raw.get("grounding", "ai")).strip().lower()
    if grounding not in ["ai", "reference", "hybrid"]:
        grounding = "ai"

    pts = raw.get("points", 0)
    try:
        pts = int(pts)
    except Exception:
        pts = 0

    if mode == "exact":
        grounding = ""

    return {
        "item_origin": "teacher_key",
        "name": str(raw.get("name", "New Item")).strip() or "New Item",
        "description": str(raw.get("description", "")).strip(),
        "points": max(0, pts),
        "grounding": grounding,
        "expected_answer": str(raw.get("expected_answer", "")).strip(),
        "mode": mode,
    }


def sanitize_item_by_origin(raw: Dict[str, Any], origin: str) -> Dict[str, Any]:
    return sanitize_teacher_key_item(raw) if origin == "teacher_key" else sanitize_assignment_item(raw)


def normalize_points(items: List[Dict[str, Any]], total_points: int, origin: str) -> List[Dict[str, Any]]:
    clean = [sanitize_item_by_origin(x, origin) for x in items if str(x.get("name", "")).strip()]
    if not clean:
        return []

    s = sum(int(x["points"]) for x in clean)
    if s <= 0:
        even = max(1, total_points // len(clean))
        clean = [{**x, "points": even} for x in clean]
        clean[0]["points"] += total_points - sum(y["points"] for y in clean)
        return clean

    out = []
    for x in clean:
        out.append({**x, "points": int(round(int(x["points"]) * total_points / s))})
    out[0]["points"] += total_points - sum(y["points"] for y in out)
    return out


def total_item_points(items: List[Dict[str, Any]]) -> int:
    return sum(int(x.get("points", 0)) for x in items)


def get_active_rubric_items_for_grade(
    grade_source: str,
    assignment_rubric_items: List[Dict[str, Any]],
    teacher_key_rubric_items: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    if grade_source == "teacher_key":
        return list(teacher_key_rubric_items or [])
    return list(assignment_rubric_items or [])


# =========================================================
# Generation
# =========================================================
def generate_items_from_assignment(assignment_text: str, total_points: int) -> Dict[str, Any]:
    payload = {
        "task": "generate_rubric_from_assignment",
        "assignment_text": assignment_text,
        "target_total_points": total_points,
        "rules": [
            "Generate grading rubric items from the assignment text.",
            "Each item must include: name, description, points, grounding.",
            "Do not include exact or conceptual mode.",
            "grounding must be ai, reference, or hybrid.",
            "Use ai when general reasoning is enough.",
            "Use reference when course material or uploaded references should be the main basis.",
            "Use hybrid when both reasoning and reference grounding are useful.",
            "Points must be integers summing exactly to target_total_points.",
        ],
        "output_json_schema": {
            "rubric_title": "string",
            "summary": ["string"],
            "items": [
                {
                    "name": "string",
                    "description": "string",
                    "points": "integer",
                    "grounding": "ai | reference | hybrid"
                }
            ]
        }
    }
    result = llm_json(payload, temperature=0.15)
    items = result.get("items", []) if isinstance(result, dict) else []
    result["items"] = normalize_points(items, total_points, "assignment")
    return result


def generate_items_from_teacher_key(
    teacher_key_text: str,
    total_points: int,
    default_grounding: Optional[str] = None,
) -> Dict[str, Any]:
    dg = (default_grounding or "").strip().lower()
    rules = [
        "Generate open-response (QA) grading criteria from the teacher's pasted or uploaded questions and model answers.",
        "Each item must include: name, description, expected_answer, points, grounding.",
        "Every criterion is graded by understanding (semantic match to the model answer), not by letter or exact string match.",
        "expected_answer should summarize the model answer or key ideas the student's response should cover.",
        "grounding must be one of: ai, reference, hybrid (never empty).",
        "Points must be integers summing exactly to target_total_points.",
    ]
    if dg in {"ai", "reference", "hybrid"}:
        rules.append(f"Use grounding '{dg}' on every item unless the teacher key clearly needs a different grounding for a specific question.")

    payload = {
        "task": "generate_rubric_from_teacher_key",
        "teacher_key_text": teacher_key_text,
        "target_total_points": total_points,
        "rules": rules,
        "output_json_schema": {
            "rubric_title": "string",
            "summary": ["string"],
            "items": [
                {
                    "name": "string",
                    "description": "string",
                    "expected_answer": "string",
                    "points": "integer",
                    "grounding": "ai | reference | hybrid",
                }
            ]
        }
    }
    result = llm_json(payload, temperature=0.1)
    items = result.get("items", []) if isinstance(result, dict) else []
    for it in items:
        if not isinstance(it, dict):
            continue
        it["mode"] = "conceptual"
        g = str(it.get("grounding", "")).strip().lower()
        if g not in {"ai", "reference", "hybrid"}:
            it["grounding"] = dg if dg in {"ai", "reference", "hybrid"} else "ai"
    result["items"] = normalize_points(items, total_points, "teacher_key")
    return result


# =========================================================
# RAG
# =========================================================
def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 120) -> List[str]:
    text = (text or "").strip()
    if not text:
        return []
    out = []
    i = 0
    while i < len(text):
        out.append(text[i:i + chunk_size])
        i += max(1, chunk_size - overlap)
    return out


def rag_retrieve_for_item(
    item: Dict[str, Any],
    ref_text: str,
    top_k: int = 3,
) -> List[str]:
    chunks = chunk_text(ref_text)
    if not chunks:
        return []

    query = " ".join([
        str(item.get("name", "")),
        str(item.get("description", "")),
        str(item.get("expected_answer", "")),
    ]).lower()

    tokens = {t for t in re.split(r"\W+", query) if len(t) > 2}
    scored = []
    for chunk in chunks:
        low = chunk.lower()
        score = sum(1 for t in tokens if t in low)
        scored.append((score, chunk))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [x[1] for x in scored[:top_k] if x[0] > 0]


def requires_reference_material(items: List[Dict[str, Any]]) -> bool:
    for item in items or []:
        grounding = str(item.get("grounding", "")).strip().lower()
        if grounding in {"reference", "hybrid"}:
            return True
    return False


def ensure_reference_documents_indexed(document_ids: Optional[List[str]] = None) -> None:
    wanted_ids = [str(document_id or "").strip() for document_id in (document_ids or []) if str(document_id or "").strip()]
    if not wanted_ids:
        return

    store = LocalVectorStore()
    indexed_ids = store.indexed_document_ids()

    for document_id in wanted_ids:
        if document_id in indexed_ids:
            continue
        document = get_local_document_by_id(document_id)
        if document is None:
            logger.warning("Reference document %s could not be resolved for indexing", document_id)
            continue
        chunks = chunk_document(document)
        if not chunks:
            logger.warning("Reference document %s produced no chunks for indexing", document_id)
            continue
        store.add_chunks(chunks)
        indexed_ids.add(document_id)


def _format_retrieved_chunk(chunk: Any) -> str:
    title = str(getattr(chunk, "document_title", "") or "").strip()
    heading = str(getattr(chunk, "section_heading", "") or "").strip()
    page = ""
    metadata = getattr(chunk, "metadata", {}) or {}
    if isinstance(metadata, dict):
        page = str(metadata.get("page_start") or metadata.get("page") or "").strip()

    prefix_parts = []
    if title:
        prefix_parts.append(title)
    if heading:
        prefix_parts.append(heading)
    if page:
        prefix_parts.append(f"page {page}")
    prefix = f"[{' | '.join(prefix_parts)}]\n" if prefix_parts else ""
    return f"{prefix}{str(getattr(chunk, 'chunk_text', '') or '').strip()}".strip()


def retrieve_reference_context_via_system_rag(
    item: Dict[str, Any],
    reference_document_ids: Optional[List[str]] = None,
    reference_text: str = "",
    top_k: int = 3,
) -> List[str]:
    contexts: List[str] = []
    query = build_reference_query(item)
    clean_document_ids = [str(document_id or "").strip() for document_id in (reference_document_ids or []) if str(document_id or "").strip()]

    if clean_document_ids and query:
        try:
            ensure_reference_documents_indexed(clean_document_ids)
            retriever = Retriever()
            rag_chunks = retriever.retrieve(query=query, top_k=top_k, document_ids=clean_document_ids)
            for chunk in rag_chunks:
                formatted = _format_retrieved_chunk(chunk)
                if formatted and formatted not in contexts:
                    contexts.append(formatted)
        except Exception:
            logger.warning("System RAG retrieval failed for grading item", exc_info=True)

    if len(contexts) < top_k and str(reference_text or "").strip():
        for snippet in retrieve_reference_context_from_text(item, reference_text, top_k=max(1, top_k - len(contexts))):
            if snippet and snippet not in contexts:
                contexts.append(snippet)

    return contexts[:top_k]


def build_reference_text_from_sources(
    reference_document_ids: Optional[List[str]] = None,
    reference_text: str = "",
) -> str:
    parts: List[str] = []

    for document_id in reference_document_ids or []:
        resolved = get_document_text(document_id)
        text = str(resolved.get("text", "")).strip()
        if text:
            parts.append(text)

    if (reference_text or "").strip():
        parts.append(str(reference_text).strip())

    return "\n\n".join(parts).strip()


def prepare_items_with_reference_context(
    items: List[Dict[str, Any]],
    reference_document_ids: Optional[List[str]] = None,
    reference_text: str = "",
) -> List[Dict[str, Any]]:
    prepared: List[Dict[str, Any]] = []
    for item in items:
        item_copy = dict(item)

        if item_copy.get("item_origin") == "assignment":
            if item_copy.get("grounding") in ["reference", "hybrid"]:
                item_copy["reference_context"] = retrieve_reference_context_via_system_rag(
                    item_copy,
                    reference_document_ids=reference_document_ids,
                    reference_text=reference_text,
                    top_k=3,
                )
            else:
                item_copy["reference_context"] = []

        elif item_copy.get("item_origin") == "teacher_key":
            if item_copy.get("mode") == "conceptual" and item_copy.get("grounding") in ["reference", "hybrid"]:
                item_copy["reference_context"] = retrieve_reference_context_via_system_rag(
                    item_copy,
                    reference_document_ids=reference_document_ids,
                    reference_text=reference_text,
                    top_k=3,
                )
            else:
                item_copy["reference_context"] = []

        else:
            item_copy["reference_context"] = []

        prepared.append(item_copy)
    return prepared


# =========================================================
# Direct reference text retrieval
# =========================================================
def build_reference_query(item: Dict[str, Any]) -> str:
    return " ".join(
        [
            str(item.get("name", "")),
            str(item.get("description", "")),
            str(item.get("expected_answer", "")),
        ]
    ).strip()


def _search_tokens(text: str) -> set[str]:
    return {tok for tok in re.findall(r"[a-z0-9]+", str(text or "").lower()) if len(tok) > 1}

def retrieve_reference_context_from_text(
    item: Dict[str, Any],
    reference_text: str,
    top_k: int = 3,
) -> List[str]:
    query = build_reference_query(item)
    source_text = str(reference_text or "").strip()
    if not query or not source_text:
        return []

    chunks = knowledge_chunk_text(source_text, chunk_size=1200, overlap=120)
    if not chunks:
        return []

    query_tokens = _search_tokens(query)
    if not query_tokens:
        return [f"[Uploaded reference] {chunks[0]}"]

    scored: List[Tuple[float, int, str]] = []
    lowered_name = str(item.get("name", "")).strip().lower()
    lowered_expected = str(item.get("expected_answer", "")).strip().lower()

    for idx, chunk in enumerate(chunks):
        chunk_tokens = _search_tokens(chunk)
        overlap = len(query_tokens & chunk_tokens)
        if overlap <= 0 and idx >= top_k:
            continue
        bonus = 0.0
        chunk_lower = chunk.lower()
        if lowered_name and lowered_name in chunk_lower:
            bonus += 2.0
        if lowered_expected and lowered_expected in chunk_lower:
            bonus += 1.5
        scored.append((overlap + bonus, idx, chunk))

    if not scored:
        return [f"[Uploaded reference] {chunks[0]}"]

    scored.sort(key=lambda row: (-row[0], row[1]))
    return [f"[Uploaded reference] {chunk}" for _, _, chunk in scored[:top_k]]


# =========================================================
# Exact-mode grading helpers (teacher_key + mode=exact)
# =========================================================

def _unique_letters(values: List[str]) -> List[str]:
    out: List[str] = []
    seen = set()
    for value in values:
        if value in seen:
            continue
        seen.add(value)
        out.append(value)
    return out


def _extract_choice_letters(text: str) -> List[str]:
    raw = html.unescape(str(text or "")).lower().strip()
    if not raw:
        return []

    matches: List[str] = []
    grouped_patterns = [
        r"\b(?:answers?|options?|choices?|selected|select|picked|marked|correct(?:\s+answers?)?)\s*(?:is|are|=|:)?\s*([a-h](?:\s*(?:,|/|&|\+|\-|and|or)\s*[a-h])*)",
        r"^\s*([a-h](?:\s*(?:,|/|&|\+|\-|and|or)\s*[a-h])*)\s*$",
    ]
    for pattern in grouped_patterns:
        for blob in re.findall(pattern, raw, flags=re.IGNORECASE):
            matches.extend(re.findall(r"\b([a-h])\b", blob, flags=re.IGNORECASE))

    numbered = re.findall(r"\b\d+\s*[-.):/]\s*([a-h])\b", raw, flags=re.IGNORECASE)
    matches.extend(numbered)

    if not matches:
        direct = re.findall(r"\b([a-h])\b", raw, flags=re.IGNORECASE)
        if 0 < len(direct) <= 5:
            matches.extend(direct)

    return _unique_letters([match.lower() for match in matches if match])


def _get_letter(text: str) -> str:
    """Extract the first MCQ option letter (a-h) from any format."""
    letters = _extract_choice_letters(text)
    if letters:
        return letters[0]

    t = text.strip().lower()
    # Explicit prefix: "A.", "A)", "(A)", "A -", "A:"
    m = re.match(r'^\(?([a-h])\)?[.):\-\s]', t)
    if m:
        return m.group(1)
    # "answer is A", "option B", "choice C"
    m = re.search(r'\b(?:answer|option|choice|select(?:ed)?)\s*(?:is\s*)?([a-h])\b', t)
    if m:
        return m.group(1)
    # Standalone letter
    m = re.fullmatch(r'([a-h])', t)
    if m:
        return m.group(1)
    return ""


def _strip_letter_prefix(text: str) -> str:
    """Remove option letter prefix from expected answer, e.g. 'B. 8' -> '8'."""
    return re.sub(r'^\(?[a-hA-H]\)?[.):\-\s]+', '', text.strip()).strip()


def _flexible_match(expected: str, student: str) -> Tuple[bool, str]:
    """
    Deterministic flexible comparison for exact-mode items.
    Handles every format: letter, content, letter+content, paraphrase, numeric, T/F.
    """
    exp = (expected or "").strip()
    stu = (student or "").strip()

    if not exp:
        return False, "No expected answer defined."
    if not stu:
        return False, "Student did not answer."

    # 1. Case-insensitive exact match
    if exp.lower() == stu.lower():
        return True, "Exact match."

    exp_letters = _extract_choice_letters(exp)
    stu_letters = _extract_choice_letters(stu)
    if exp_letters and stu_letters:
        exp_set = set(exp_letters)
        stu_set = set(stu_letters)
        if exp_set == stu_set:
            label = ", ".join(letter.upper() for letter in exp_letters)
            if len(exp_set) > 1:
                return True, f"Multiple-choice set match: {label}."
            return True, f"Option letter match: {label}."
        if len(exp_set) > 1 or len(stu_set) > 1:
            exp_label = ", ".join(letter.upper() for letter in exp_letters)
            stu_label = ", ".join(letter.upper() for letter in stu_letters)
            return False, f"Option set mismatch: expected {exp_label}, got {stu_label}."

    exp_letter = _get_letter(exp)
    stu_letter = _get_letter(stu)
    exp_content = _strip_letter_prefix(exp)      # e.g. "B. 8" -> "8"
    stu_content = _strip_letter_prefix(stu)      # strip letter from student too

    # 2. Both have letters and they match
    if exp_letter and stu_letter and exp_letter == stu_letter:
        return True, f"Option letter match: {exp_letter.upper()}."

    # 3. Student wrote the content of the correct option (e.g. "8" when expected is "B. 8")
    if exp_content:
        ec = exp_content.lower()
        sc = stu.lower()
        sc2 = stu_content.lower()
        if ec == sc or ec == sc2:
            return True, f"Content match: '{exp_content}'."
        # Content contained within student answer
        if ec in sc or sc in ec:
            return True, f"Content contained in answer."
        # Numeric match (e.g. 8 == 8.0)
        if numeric_match(exp_content, stu):
            return True, "Numeric match."

    # 4. True/False
    exp_tf = extract_true_false(exp)
    stu_tf = extract_true_false(stu)
    if exp_tf and stu_tf:
        return (exp_tf == stu_tf), ("True/False match." if exp_tf == stu_tf else f"True/False mismatch: expected {exp_tf}, got {stu_tf}.")

    # 5. If expected is just a letter, also check student content against that letter
    if exp_letter and not exp_content:
        if stu_letter == exp_letter:
            return True, f"Letter match: {exp_letter.upper()}."

    # 6. Token similarity on content
    cmp_exp = exp_content if exp_content else exp
    sim = token_similarity(cmp_exp, stu)
    if sim >= 0.82:
        return True, f"High token similarity ({sim:.0%})."

    return False, f"No match (similarity {sim:.0%})."


def _is_hard_exact_mismatch(expected: str, student: str) -> Tuple[bool, str]:
    exp_letters = _extract_choice_letters(expected)
    stu_letters = _extract_choice_letters(student)
    if exp_letters and stu_letters:
        exp_set = set(exp_letters)
        stu_set = set(stu_letters)
        if exp_set != stu_set:
            exp_label = ", ".join(letter.upper() for letter in exp_letters)
            stu_label = ", ".join(letter.upper() for letter in stu_letters)
            if len(exp_set) > 1 or len(stu_set) > 1:
                return True, f"Option set mismatch: expected {exp_label}, got {stu_label}."
            return True, f"Option letter mismatch: expected {exp_label}, got {stu_label}."

    exp_tf = extract_true_false(expected)
    stu_tf = extract_true_false(student)
    if exp_tf and stu_tf and exp_tf != stu_tf:
        return True, f"True/False mismatch: expected {exp_tf}, got {stu_tf}."

    return False, ""


def normalize_answer_text(text: str) -> str:
    text = html.unescape(str(text or "")).lower().strip()
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"[\(\)\[\]\{\}:;,.!?]", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_mcq_choice(text: str) -> str:
    raw = str(text or "").strip()
    t = normalize_answer_text(raw)
    explicit_patterns = [
        r"\banswer\s*(?:is|=)?\s*([a-h])\b",
        r"\bcorrect\s*(?:answer)?\s*(?:is|=)?\s*([a-h])\b",
        r"\bselected\s*(?:answer)?\s*(?:is|=)?\s*([a-h])\b",
        r"\bchosen\s*(?:answer)?\s*(?:is|=)?\s*([a-h])\b",
        r"\bpicked\s*(?:answer)?\s*(?:is|=)?\s*([a-h])\b",
        r"\bmarked\s*(?:answer)?\s*(?:is|=)?\s*([a-h])\b",
        r"\bchoice\s*([a-h])\b",
        r"\boption\s*([a-h])\b",
    ]
    for pattern in explicit_patterns:
        m = re.search(pattern, t)
        if m:
            return m.group(1).lower()

    option_list_matches = re.findall(r"(?im)^\s*([a-h])[\).:-]\s+\S", raw)
    lines = [normalize_answer_text(line) for line in raw.splitlines() if line.strip()]

    if len(option_list_matches) >= 2:
        for line in reversed(lines[-3:]):
            m = re.fullmatch(r"([a-h])", line)
            if m:
                return m.group(1).lower()
        return ""

    for line in reversed(lines[-3:]):
        m = re.fullmatch(r"([a-h])", line)
        if m:
            return m.group(1).lower()
        m = re.fullmatch(r"([a-h])\s+\S.*", line)
        if m and not re.search(r"\b[a-h]\s+\S+\s+[a-h]\s+\S+", line):
            return m.group(1).lower()

    fallback_patterns = [
        r"^([a-h])$",
        r"\b([a-h])\b",
    ]
    for pattern in fallback_patterns:
        m = re.search(pattern, t)
        if m:
            return m.group(1).lower()
    return ""


def extract_true_false(text: str) -> str:
    tokens = set(normalize_answer_text(text).split())
    if tokens & {"true", "t", "yes", "correct"}:
        return "true"
    if tokens & {"false", "f", "no", "incorrect"}:
        return "false"
    return ""


def numeric_match(expected: str, student: str, tolerance: float = 1e-6) -> bool:
    exp_nums = re.findall(r"[-+]?\d+(?:\.\d+)?", str(expected or ""))
    stu_nums = re.findall(r"[-+]?\d+(?:\.\d+)?", str(student or ""))
    if not exp_nums or not stu_nums:
        return False
    try:
        expected_values = [float(value) for value in exp_nums]
        student_values = [float(value) for value in stu_nums]
        return all(
            any(abs(expected_value - student_value) <= tolerance for student_value in student_values)
            for expected_value in expected_values
        )
    except Exception:
        return False


def token_similarity(a: str, b: str) -> float:
    stop = {"the", "a", "an", "and", "or", "to", "of", "in", "on", "for", "is", "are", "answer", "choice", "option"}
    ta = {x for x in normalize_answer_text(a).split() if x not in stop and (len(x) > 1 or x.isdigit())}
    tb = {x for x in normalize_answer_text(b).split() if x not in stop and (len(x) > 1 or x.isdigit())}
    if not ta or not tb:
        return 0.0
    return len(ta & tb) / len(ta | tb)


def exact_answers_match(expected: str, student: str) -> Tuple[bool, str]:
    exp = normalize_answer_text(expected)
    stu = normalize_answer_text(student)

    if not exp or not stu:
        return False, "Missing expected answer or student answer."

    if exp == stu:
        return True, "Exact text match."

    exp_tf = extract_true_false(expected)
    stu_tf = extract_true_false(student)
    if exp_tf and stu_tf:
        return (exp_tf == stu_tf, "True/false match." if exp_tf == stu_tf else "True/false mismatch.")

    exp_choice = extract_mcq_choice(expected)
    stu_choice = extract_mcq_choice(student)
    if exp_choice and stu_choice:
        return (exp_choice == stu_choice, "MCQ option match." if exp_choice == stu_choice else "MCQ option mismatch.")

    if numeric_match(expected, student):
        return True, "Numeric match."

    if exp in stu or stu in exp:
        return True, "Answer text is contained in the other answer."

    sim = token_similarity(expected, student)
    if sim >= 0.85:
        return True, f"Near-exact token overlap ({sim:.2f})."

    return False, f"No exact match (similarity {sim:.2f})."


def _reduce_answer_chunk(chunk: str) -> str:
    text = str(chunk or "").strip()
    if not text:
        return ""

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    explicit_patterns = [
        r"(?i)\b(?:student\s*)?answer\s*[:=\-]",
        r"(?i)\bselected\s*(?:answer)?\s*[:=\-]",
        r"(?i)\bchosen\s*(?:answer)?\s*[:=\-]",
        r"(?i)\bpicked\s*(?:answer)?\s*[:=\-]",
        r"(?i)\bmarked\s*(?:answer)?\s*[:=\-]",
    ]
    for line in reversed(lines):
        if any(re.search(pattern, line) for pattern in explicit_patterns):
            return line

    return text


def _extract_inline_answer_fragment(line: str, question_number: str = "") -> str:
    raw = str(line or "").strip()
    if not raw:
        return ""

    patterns = [
        r"(?i)\b(?:student\s*)?answer\s*[:=\-]\s*(.+)$",
        r"(?i)\bans\s*[:=\-]\s*(.+)$",
        r"(?i)\b(?:selected|chosen|picked|marked)\s*(?:answer)?\s*[:=\-]?\s*(.+)$",
    ]
    if question_number:
        qn = re.escape(question_number)
        patterns.extend(
            [
                rf"(?i)^(?:question\s*|q\s*)?{qn}\s*[-.):/]\s*(.+)$",
                rf"(?i)\b(?:question\s*|q\s*)?{qn}\s*[-.):/]\s*([a-h](?:\s*(?:,|/|&|\+|\-|and|or)\s*[a-h])+|[a-h]|true|false|yes|no|\d+(?:[.,]\d+)?|.+)$",
                rf"(?i)\b{qn}\s*-\s*([a-h](?:\s*(?:,|/|&|\+|\-|and|or)\s*[a-h])+|[a-h]|.+)$",
            ]
        )

    for pattern in patterns:
        match = re.search(pattern, raw)
        if not match:
            continue
        candidate = match.group(1).strip(" -:\t")
        if candidate:
            return candidate
    return ""


def _looks_like_answer_candidate(line: str) -> bool:
    raw = str(line or "").strip()
    if not raw:
        return False
    if raw.endswith("?"):
        return False
    if len(raw.split()) > 14:
        return False
    if re.match(r"(?i)^(question|q)\s*\d+\b", raw):
        return False
    return True


def _extract_answer_from_nearby_lines(chunk: str, question_number: str = "") -> str:
    lines = [line.strip() for line in str(chunk or "").splitlines() if line.strip()]
    if not lines:
        return ""

    for line in reversed(lines):
        candidate = _extract_inline_answer_fragment(line, question_number)
        if candidate:
            return candidate

    if question_number:
        joined = " \n ".join(lines)
        numbered = re.search(
            rf"(?i)\b(?:question\s*|q\s*)?{re.escape(question_number)}\s*[-.):/]\s*([^\n]+)",
            joined,
        )
        if numbered:
            candidate = numbered.group(1).strip(" -:\t")
            if candidate:
                return candidate

    for line in reversed(lines[-4:]):
        if _looks_like_answer_candidate(line):
            return line

    return ""


def extract_student_answer_for_item(submission_text: str, item: Dict[str, Any], index: int = 0) -> str:
    text = str(submission_text or "").strip()
    if not text:
        return ""

    name = str(item.get("name", ""))
    desc = str(item.get("description", ""))
    candidates = [name, desc]
    question_number = ""

    m = re.search(r"(?:question|q)\s*(\d+)", name, flags=re.I)
    if m:
        question_number = m.group(1)
        candidates += [
            f"Question {question_number}",
            f"Q{question_number}",
            f"Q {question_number}",
            f"{question_number}.",
            f"{question_number})",
            f"{question_number}-",
            f"{question_number}:",
        ]
    else:
        nums = re.findall(r"\d+", name)
        if nums:
            question_number = nums[0]
            candidates += [
                f"Question {question_number}",
                f"Q{question_number}",
                f"Q {question_number}",
                f"{question_number}.",
                f"{question_number})",
                f"{question_number}-",
                f"{question_number}:",
            ]

    lower = text.lower()
    for cand in candidates:
        cand = cand.strip()
        if not cand:
            continue
        pos = lower.find(cand.lower())
        if pos >= 0:
            start = pos + len(cand)
            chunk = text[start:start + 1500].strip()
            nxt = re.search(r"\n\s*(?:question\s+\d+|q\s*\d+|\d+[\).:-])", chunk, flags=re.I)
            if nxt and nxt.start() > 5:
                chunk = chunk[:nxt.start()]
            chunk = re.sub(r"^\s*[:\-â€“â€”]?\s*", "", chunk).strip()
            chunk = re.sub(r"^\s*(student\s*)?answer\s*[:\-â€“â€”]\s*", "", chunk, flags=re.I).strip()
            if chunk:
                candidate = _extract_answer_from_nearby_lines(chunk, question_number)
                if candidate:
                    return candidate
                return _reduce_answer_chunk(chunk)

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if index < len(lines):
        line = lines[index]
        line = re.sub(r"^\s*(?:\d+|q\d+)\s*[\).:-]\s*", "", line, flags=re.I).strip()
        candidate = _extract_answer_from_nearby_lines(line, question_number)
        if candidate:
            return candidate
        return _reduce_answer_chunk(line)

    if len(text) <= 500:
        candidate = _extract_answer_from_nearby_lines(text, question_number)
        if candidate:
            return candidate
        return _reduce_answer_chunk(text)

    return ""


def make_exact_result(item: Dict[str, Any], earned: float, rationale: str, student_answer: str) -> Dict[str, Any]:
    points = _coerce_int_points(item.get("points"))
    return {
        "item_origin": "teacher_key",
        "name": item.get("name", ""),
        "description": item.get("description", ""),
        "expected_answer": item.get("expected_answer", ""),
        "points": points,
        "mode": "exact",
        "grounding": "",
        "earned_points": round(float(earned), 2),
        "rationale": rationale,
        "suggestions": [] if earned == points else [f"Correct answer: {item.get('expected_answer', '')}"],
        "evidence": [{"quote": student_answer[:500], "source": "student_submission"}] if student_answer else [],
    }


def _regex_extract_answers(submission_text: str, items: List[Dict[str, Any]]) -> Dict[str, str]:
    """Conservative regex extraction — only returns an answer when it is CERTAIN.
    Better to return nothing and fall back to LLM than to return the wrong thing.

    Matches ONLY these safe patterns:
      • "1. A"  "1) B"  "1: C"  — whole line is number + single letter
      • "Q1 A"  "Q1. B"         — Q-prefix + number + single letter
      • "1. True" / "1. False" / "1. Yes" / "1. No"
      • "1. 8"  "1. 42"         — whole line is number + bare number
      • Explicit "Answer: X" or "Ans: X" anywhere on a line near the question
    """
    results: Dict[str, str] = {}
    lines = [l.strip() for l in submission_text.splitlines() if l.strip()]
    multi_choice_blob = r"[a-h](?:\s*(?:,|/|&|\+|\-|and|or)\s*[a-h])+"
    answer_blob = rf"(?:{multi_choice_blob}|[a-h]|true|false|yes|no|\d+(?:[.,]\d+)?)"

    # Build a quick index: line_index → which question number(s) appear on it
    # so we can look for "Answer:" on the line(s) after a question
    for item in items:
        name = item.get("name", "")
        nums = re.findall(r'\d+', name)
        if not nums:
            continue
        n = nums[0]

        for idx, line in enumerate(lines):
            # ── Pattern A: entire line is  [Q|Question] N [sep] SINGLE_ANSWER ──
            # SINGLE_ANSWER = one letter, true/false, yes/no, or a bare number
            m = re.match(
                r'^(?:q(?:uestion)?\s*)?' + re.escape(n) +
                r'\s*[.):\-\s]\s*'
                + rf'({answer_blob})'
                r'\s*$',
                line, re.IGNORECASE
            )
            if m:
                results[name] = m.group(1).strip()
                break

            # ── Pattern B: explicit "Answer: X" or "Ans: X" on this line ──
            # Only accept if the IMMEDIATELY preceding line starts with question N
            ans_m = re.search(r'\bans(?:wer)?\s*[:\-]\s*(.+)', line, re.IGNORECASE)
            if ans_m and idx > 0:
                prev = lines[idx - 1]
                # Previous line must start with the question number (as a question marker)
                if re.match(r'^(?:q(?:uestion)?\s*)?' + re.escape(n) + r'[.):\s]', prev, re.IGNORECASE):
                    raw = ans_m.group(1).strip()
                    if raw and not raw.endswith('?') and len(raw.split()) <= 10:
                        results[name] = raw
                        break

    return results


def _extract_student_answers_llm(
    submission_text: str,
    items: List[Dict[str, Any]],
) -> Dict[str, str]:
    """Extract student answers: regex first, then local heuristics, then LLM."""

    # Pass 1: fast regex for standard MCQ formats
    regex_results = _regex_extract_answers(submission_text, items)
    heuristic_results: Dict[str, str] = {}

    for index, item in enumerate(items):
        name = str(item.get("name", "")).strip()
        if not name or _fuzzy_name_lookup(regex_results, name):
            continue
        candidate = extract_student_answer_for_item(submission_text, item, index=index)
        if candidate and len(candidate.split()) <= 16:
            heuristic_results[name] = candidate

    extracted_results = dict(regex_results)
    extracted_results.update(heuristic_results)

    missing = [i for i in items if not _fuzzy_name_lookup(extracted_results, i.get("name", ""))]
    if not missing:
        logger.debug("All %d answers found via regex — skipping LLM extraction", len(items))
        return extracted_results

    logger.debug(
        "Deterministic extraction found %d/%d; LLM will handle %d remaining",
        len(extracted_results),
        len(items),
        len(missing),
    )

    # Pass 2: LLM for remaining items
    payload = {
        "task": "extract_student_answers_from_submission",
        "student_submission": submission_text,
        "questions": [
            {
                "name": i.get("name", ""),
                "question_text": i.get("description", ""),
                "candidate_answer": extract_student_answer_for_item(submission_text, i),
            }
            for i in missing
        ],
        "rules": [
            "Read the ENTIRE student submission carefully.",
            "For each question, find what the student wrote as their answer.",
            "Match questions by number ('1.', 'Q1', 'Question 1', '(1)', '1)' all mean question 1).",
            "candidate_answer is a local guess from nearby text in the submission. Use it as a hint, but verify it against the actual submission text.",
            "Extract ONLY what the student wrote — a letter, multiple option letters, a word, a number, or a short phrase.",
            "If a question has multiple selected options, preserve all of them, for example 'A, C' or 'B and D'.",
            "For PDFs, the answer may appear near the question in forms like '4-B', 'Q4: B', 'Answer: B', or as a short word/phrase on the next line.",
            "Do NOT include the question text itself in the answer.",
            "If a question has no answer at all, set student_answer to empty string.",
        ],
        "output_json_schema": {
            "answers": [
                {"name": "string", "student_answer": "string"}
            ]
        },
    }
    try:
        res = llm_json(payload, temperature=0.0)
        llm_answers = res.get("answers", []) if isinstance(res, dict) else []
        for a in llm_answers:
            k = str(a.get("name", "")).strip()
            v = str(a.get("student_answer", "")).strip()
            if k and v:
                extracted_results[k] = v
    except Exception:
        logger.warning("LLM answer extraction failed", exc_info=True)

    return extracted_results


def _fuzzy_name_lookup(extracted: Dict[str, str], target_name: str) -> str:
    """Try exact match first, then case-insensitive, then check if target
    number appears in any extracted key (handles 'Q1' vs 'Question 1')."""
    if target_name in extracted:
        return extracted[target_name]
    lower = target_name.lower()
    for key, val in extracted.items():
        if key.lower() == lower:
            return val
    # Try matching by the numeric part (e.g. "1" matches "Question 1", "Q1", "1.")
    nums = re.findall(r"\d+", target_name)
    if nums:
        n = nums[0]
        for key, val in extracted.items():
            key_nums = re.findall(r"\d+", key)
            if key_nums and key_nums[0] == n:
                return val
    return ""


def _grade_extracted_answers_llm(
    items_with_answers: List[Dict[str, Any]],
    teacher_key_text: str,
    submission_text: str,
) -> List[Dict[str, Any]]:
    """Grade items that have pre-extracted student answers.
    The full submission is also included so the LLM can recover
    if extraction missed an answer."""
    payload = {
        "task": "grade_exact_items_flexibly",
        "teacher_key": teacher_key_text,
        "student_submission": submission_text,
        "grading_items": [
            {
                "name": i.get("name"),
                "description": i.get("description", ""),
                "expected_answer": i.get("expected_answer", ""),
                "points": i.get("points", 0),
                "student_answer": i.get("_student_answer", ""),
            }
            for i in items_with_answers
        ],
        "rules": [
            "student_answer already contains what the student wrote if extraction succeeded.",
            "If student_answer is empty, search the student_submission yourself to find the answer for that question.",
            "CRITICAL — expected_answer format handling:",
            "  If expected_answer contains multiple correct options like 'A, C' or 'B and D':",
            "    → The student must provide the full correct set of options and no wrong extras.",
            "    → Accept separators like commas, slashes, ampersands, hyphens, or the word 'and'.",
            "    → Accept formats like '1-A, 3-C', 'A and C', 'A/C', or 'A+C'.",
            "  If expected_answer looks like 'B. 8' or 'C. photosynthesis' (letter + dot + content):",
            "    → The LETTER part ('B', 'C') and the CONTENT part ('8', 'photosynthesis') are BOTH valid answers.",
            "    → If the student wrote ONLY the content ('8', 'photosynthesis') — CORRECT.",
            "    → If the student wrote ONLY the letter ('B', 'C') — CORRECT.",
            "    → If the student wrote both ('B. 8') — CORRECT.",
            "    → If the student wrote a paraphrase of the content — CORRECT.",
            "  If expected_answer is just a letter like 'B':",
            "    → Look up option B in the teacher_key to find its text.",
            "    → Accept the letter OR the option text OR a paraphrase of it.",
            "  If expected_answer is just content text (no letter prefix):",
            "    → Accept exact match, paraphrase, synonyms, abbreviations.",
            "    → Also accept if student wrote the corresponding option letter.",
            "ALL of these student answer formats count as correct:",
            "  • Letter alone: 'A', 'b', 'C'",
            "  • Letter with punctuation: 'A)', '(A)', 'A.'",
            "  • Content/value alone: '8', 'photosynthesis', 'the nervous system'",
            "  • Letter + content: 'B. 8', 'C - photosynthesis'",
            "  • Paraphrase: 'eight', 'plants making food', 'nervous system controls body'",
            "  • True/False synonyms: T, F, Yes, No, Correct, Incorrect",
            "  • Numeric: accept rounding differences (e.g. 8.0 = 8)",
            "Scoring is BINARY: full points if correct in any format, 0 if genuinely wrong.",
            "ONLY award 0 if the student's answer is clearly a different/wrong answer, not just a different format.",
        ],
        "output_json_schema": {
            "items_results": [
                {
                    "name": "string",
                    "earned_points": "number (full points or 0)",
                    "rationale": "string",
                    "student_answer": "string — what the student wrote"
                }
            ]
        },
    }
    try:
        res = llm_json(payload, temperature=0.0)
        return res.get("items_results", []) if isinstance(res, dict) else []
    except Exception:
        logger.warning("LLM exact grading failed", exc_info=True)
        return []


def _grade_exact_items_via_llm(
    submission_text: str,
    teacher_key_text: str,
    exact_items: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    # Step 1: LLM extracts student answers (what the student wrote for each question)
    extracted = _extract_student_answers_llm(submission_text, exact_items)

    out: List[Dict[str, Any]] = []
    uncertain: List[Dict[str, Any]] = []

    for item in exact_items:
        name = item.get("name", "")
        points = int(item.get("points", 0))
        expected = str(item.get("expected_answer", "")).strip()
        student_answer = _fuzzy_name_lookup(extracted, name)

        if not student_answer:
            # Extraction missed — defer to LLM
            uncertain.append({**item, "_student_answer": ""})
            continue

        # Step 2: Deterministic comparison
        matched, reason = _flexible_match(expected, student_answer)
        if matched:
            out.append({
                "item_origin": "teacher_key",
                "name": name,
                "description": item.get("description", ""),
                "expected_answer": expected,
                "points": points,
                "mode": "exact",
                "grounding": "",
                "earned_points": float(points),
                "rationale": f"Student wrote: '{student_answer}'. {reason}",
                "suggestions": [],
                "evidence": [{"quote": student_answer[:300], "source": "student_submission"}],
            })
        else:
            hard_mismatch, hard_reason = _is_hard_exact_mismatch(expected, student_answer)
            if hard_mismatch:
                out.append(make_exact_result(item, 0.0, f"Student wrote: '{student_answer}'. {hard_reason}", student_answer))
                continue
            # Deterministic missed — let LLM check semantic equivalence
            # (e.g. "eight" == "8", synonyms, paraphrases)
            uncertain.append({**item, "_student_answer": student_answer})

    # Step 3: LLM fallback for items where deterministic matching couldn't decide
    if uncertain:
        uncertain_by_name = {i.get("name", ""): i for i in uncertain}
        fallback = _grade_extracted_answers_llm(uncertain, teacher_key_text, submission_text)
        fallback_by_name = {str(r.get("name", "")).strip(): r for r in fallback}

        for item in uncertain:
            name = item.get("name", "")
            r = fallback_by_name.get(name, {})
            stu_ans = str(r.get("student_answer", "") or item.get("_student_answer", "")).strip()
            expected = str(item.get("expected_answer", "")).strip()
            pts = int(item.get("points", 0))

            if stu_ans:
                # Re-apply deterministic check on what the LLM found/confirmed
                matched, reason = _flexible_match(expected, stu_ans)
                if not matched:
                    hard_mismatch, hard_reason = _is_hard_exact_mismatch(expected, stu_ans)
                    if hard_mismatch:
                        earned = 0.0
                        rationale = f"Student wrote: '{stu_ans}'. {hard_reason}"
                        out.append(make_exact_result(item, earned, rationale, stu_ans))
                        continue
                    # Trust LLM's binary verdict as last resort
                    try:
                        llm_earned = float(r.get("earned_points", 0))
                        matched = llm_earned > 0 and pts > 0
                    except Exception:
                        matched = False
                earned = float(pts) if matched else 0.0
                if matched and float(earned) > 0 and reason.startswith("No match"):
                    rationale = f"Student wrote: '{stu_ans}'. {str(r.get('rationale', 'Accepted by fallback semantic check.')).strip()}"
                else:
                    rationale = f"Student wrote: '{stu_ans}'. {reason if matched else r.get('rationale', 'No match found.')}"
            else:
                earned = 0.0
                rationale = "Answer not found in submission."

            # Always use make_exact_result so every field is properly populated
            out.append(make_exact_result(item, earned, rationale, stu_ans))

    by_name = {r["name"]: r for r in out}
    return [
        by_name.get(str(item.get("name", "") or "").strip())
        or make_exact_result(item, 0.0, "Answer not found in submission.", "")
        for item in exact_items
    ]

# =========================================================
# Fast grading
# =========================================================
def _max_chars_grade_prompt() -> int:
    """Upper bound on student + teacher key text in one Groq grading call (TPM on on_demand tier)."""
    raw = (os.getenv("GROQ_GRADE_PROMPT_MAX_CHARS") or "").strip()
    if raw.isdigit():
        return max(4_000, min(int(raw), 500_000))
    return 12_000


def _grade_items_per_llm_call() -> int | None:
    """Rubric rows per Groq request. ``None`` = entire non-exact rubric in one call (fastest, most consistent).

    Set ``GRADE_ITEMS_PER_LLM_CALL`` to a positive integer (e.g. 3) to split very large rubrics and reduce TPM spikes.
    ``0`` or ``all`` means one call for the whole rubric.
    """
    raw = (os.getenv("GRADE_ITEMS_PER_LLM_CALL") or "0").strip().lower()
    if raw in ("0", "", "all", "none"):
        return None
    if raw.isdigit():
        v = int(raw)
        if v <= 0:
            return None
        return max(1, min(v, 200))
    return None


def _grade_chunk_gap_seconds() -> float:
    """Pause between chunked grading calls (only when ``GRADE_ITEMS_PER_LLM_CALL`` > 0)."""
    raw = (os.getenv("GRADE_LLM_CHUNK_GAP_SEC") or "0").strip()
    try:
        return max(0.0, min(float(raw), 60.0))
    except ValueError:
        return 0.0


def _truncate_grade_text(text: str, label: str) -> str:
    t = (text or "").strip()
    cap = _max_chars_grade_prompt()
    if len(t) <= cap:
        return t
    head = t[:cap].rstrip()
    last_break = max(head.rfind("\n\n"), head.rfind(". "))
    if last_break > int(cap * 0.5):
        head = head[:last_break].rstrip()
    note = (
        f"\n\n[… {label} truncated for model limits ({cap} chars); "
        "grade only from the text above.]\n"
    )
    return head + note


def build_grade_failure_result(
    submission_text: str,
    items: List[Dict[str, Any]],
    error_message: str,
) -> Dict[str, Any]:
    """Deterministic result when grading cannot finish (batch resilience, UI feedback)."""
    msg = (error_message or "Grading failed.").strip()[:800]
    final_results: List[Dict[str, Any]] = []
    total_earned = 0.0
    total_possible = 0

    for item in items:
        if item.get("item_origin") == "teacher_key" and item.get("mode") == "exact":
            r = make_exact_result(item, 0.0, msg, "")
            final_results.append(r)
            total_earned += float(r.get("earned_points", 0) or 0)
            total_possible += _coerce_int_points(item.get("points"))
            continue

        pts = _coerce_int_points(item.get("points"))
        total_possible += pts
        final_results.append({
            "item_origin": item.get("item_origin", ""),
            "name": item.get("name", ""),
            "description": item.get("description", ""),
            "expected_answer": item.get("expected_answer", ""),
            "points": pts,
            "mode": item.get("mode", ""),
            "grounding": item.get("grounding", ""),
            "earned_points": 0.0,
            "rationale": msg,
            "suggestions": ["Retry grading after checking the API key, model, and network."],
            "evidence": [],
            "matched_key_ideas": [],
            "missing_key_ideas": [],
            "misconceptions": [],
        })

    return {
        "overall_score": round(total_earned, 2),
        "overall_out_of": total_possible,
        "items_results": final_results,
    }


def grade_submission_fast(
    submission_text: str,
    items: List[Dict[str, Any]],
    teacher_key_text: str,
    reference_document_ids: Optional[List[str]] = None,
    reference_text: str = "",
    batch_submission: bool = False,
) -> Dict[str, Any]:
    exact_items = [
        item for item in items
        if item.get("item_origin") == "teacher_key" and item.get("mode") == "exact"
    ]
    non_exact_items = [
        item for item in items
        if not (item.get("item_origin") == "teacher_key" and item.get("mode") == "exact")
    ]

    by_name: Dict[str, Any] = {}

    if exact_items:
        for result in _grade_exact_items_via_llm(submission_text, teacher_key_text, exact_items):
            name = str(result.get("name", "")).strip()
            if name:
                by_name[name] = result

    if non_exact_items:
        prepared_full = prepare_items_with_reference_context(
            non_exact_items,
            reference_document_ids=reference_document_ids,
            reference_text=reference_text,
        )
        stu = _truncate_grade_text(submission_text, "Student submission")
        tkey = _truncate_grade_text(teacher_key_text, "Teacher key")
        if batch_submission:
            # Adaptive: one Groq call when the rubric is small (fast). Larger rubrics use fewer, bigger chunks
            # than before (was 2 rows + long gaps) so batch stays TPM-safe without taking many minutes.
            n = len(prepared_full)
            try:
                single_max = max(1, min(int((os.getenv("GRADE_BATCH_SINGLE_CALL_MAX_ITEMS") or "10").strip() or "10"), 80))
            except ValueError:
                single_max = 10
            override = (os.getenv("GRADE_BATCH_ITEMS_PER_CALL") or "").strip()
            if override.isdigit() and int(override) > 0:
                chunk_n = max(1, min(int(override), 80))
                try:
                    gap = max(0.0, float((os.getenv("GRADE_BATCH_CHUNK_GAP_SEC") or "0.35").strip() or "0.35"))
                except ValueError:
                    gap = 0.35
            elif n <= single_max:
                chunk_n = n
                gap = 0.0
            else:
                chunk_n = max(4, min(8, (n + 2) // 3))
                try:
                    gap = max(0.0, float((os.getenv("GRADE_BATCH_CHUNK_GAP_SEC") or "0.35").strip() or "0.35"))
                except ValueError:
                    gap = 0.35
        else:
            chunk_n = _grade_items_per_llm_call()
            gap = _grade_chunk_gap_seconds()
        rules = [
            "Assignment-origin items do not use exact/conceptual mode. Grade them using their grounding only.",
            "If an assignment-origin item grounding is ai, use reasoning.",
            "If an assignment-origin item grounding is reference, prioritize reference_context.",
            "If an assignment-origin item grounding is hybrid, use both reasoning and reference_context.",
            "Teacher-key-origin conceptual items should grade by understanding, not exact wording.",
            "For teacher-key-origin conceptual items: use grounding ai/reference/hybrid.",
            "Return one result per rubric row in grading_items with earned_points, rationale, suggestions, evidence.",
            "earned_points must be between 0 and that item's points.",
            "Return overall_score and overall_out_of for the rows in grading_items only (subset totals).",
            "grading_items may be a subset of the full assignment rubric; only grade those rows.",
            "ASSIGNMENT-ORIGIN ONLY (item_origin is assignment): be a skeptical grader. Default to the lower defensible score when evidence is thin; do not inflate.",
            "ASSIGNMENT-ORIGIN ONLY: full points only for exceptional work that explicitly and completely satisfies the criterion with no gaps, vagueness, or missing sub-parts implied by the description.",
            "ASSIGNMENT-ORIGIN ONLY: if the criterion implies multiple distinct elements, require clear evidence for each; partial coverage earns proportional credit in the lower half of the points band, not 'nearly full' scores.",
            "ASSIGNMENT-ORIGIN ONLY: penalize padded length, generic platitudes, restating the assignment prompt, and topic-adjacent filler that does not directly prove the criterion.",
            "ASSIGNMENT-ORIGIN ONLY: when awarding more than half of a row's points, quote or paraphrase concrete student wording in rationale or evidence for every such row.",
            "ASSIGNMENT-ORIGIN ONLY: 'mostly right' or broadly on-topic should typically fall around 35–65% of that row's points, not 85–100%, unless the text is truly precise and complete.",
            "TEACHER-KEY CONCEPTUAL ONLY (item_origin is teacher_key, mode is conceptual): keep the usual fair semantic standard — equivalent paraphrases and correct reasoning still earn credit; do not apply the assignment-only strictness above to these rows.",
        ]
        schema = {
            "overall_score": "number",
            "overall_out_of": "number",
            "items_results": [
                {
                    "name": "string",
                    "earned_points": "number",
                    "rationale": "string",
                    "suggestions": ["string"],
                    "evidence": [{"quote": "string", "source": "string"}],
                }
            ],
        }
        if chunk_n is None or chunk_n >= len(prepared_full):
            slices = [prepared_full]
        else:
            slices = [prepared_full[i : i + chunk_n] for i in range(0, len(prepared_full), chunk_n)]

        for idx, chunk in enumerate(slices):
            if idx > 0 and gap > 0:
                time.sleep(gap)
            payload = {
                "task": "grade_entire_submission",
                "student_submission": stu,
                "teacher_key_text": tkey,
                "grading_items": chunk,
                "rules": rules,
                "output_json_schema": schema,
            }
            # temperature=0 and optional seed (Groq) improve repeatability; assignment strictness is rule-driven.
            res = llm_json(payload, temperature=0.0)
            raw_results = res.get("items_results", []) if isinstance(res, dict) else []
            for r in raw_results:
                name = str(r.get("name", "")).strip()
                if name:
                    by_name[name] = r

    final_results = []
    total_earned = 0.0
    total_possible = 0

    for item in items:
        item_name = str(item.get("name", "") or "").strip()
        raw = by_name.get(item_name, {}) if item_name else {}
        if item.get("item_origin") == "teacher_key" and item.get("mode") == "exact":
            result = raw if raw else make_exact_result(item, 0.0, "Answer not found in submission.", "")
            earned = float(result.get("earned_points", 0) or 0)
            total_earned += earned
            total_possible += _coerce_int_points(item.get("points"))
            final_results.append(result)
            continue

        try:
            earned = float(raw.get("earned_points", 0))
        except Exception:
            earned = 0.0
        row_max = float(_coerce_int_points(item.get("points")))
        earned = max(0.0, min(row_max, earned))
        total_earned += earned
        total_possible += _coerce_int_points(item.get("points"))

        final_results.append({
            "item_origin": item.get("item_origin", ""),
            "name": item.get("name", ""),
            "description": item.get("description", ""),
            "expected_answer": item.get("expected_answer", ""),
            "points": item.get("points", 0),
            "mode": item.get("mode", ""),
            "grounding": item.get("grounding", ""),
            "earned_points": round(earned, 2),
            "rationale": raw.get("rationale", ""),
            "suggestions": raw.get("suggestions", []),
            "evidence": raw.get("evidence", []),
            "matched_key_ideas": raw.get("matched_key_ideas", []),
            "missing_key_ideas": raw.get("missing_key_ideas", []),
            "misconceptions": raw.get("misconceptions", []),
        })

    return {
        "overall_score": round(total_earned, 2),
        "overall_out_of": total_possible,
        "items_results": final_results,
    }

def build_result_record(
    title: str,
    result: Dict[str, Any],
    submission_text: str,
    history_type: str = "single",
    batch_id: str = "",
    batch_name: str = "",
    batch_size: int = 0,
    batch_rank: int = 0,
    batch_created_at: str = "",
) -> Dict[str, Any]:
    return normalize_result_record({
        "id": f"{datetime.now().strftime('%Y%m%d-%H%M%S')}-{short_hash(submission_text, 8)}",
        "timestamp": now_iso(),
        "title": title,
        "overall_score": result.get("overall_score", 0),
        "overall_out_of": result.get("overall_out_of", 0),
        "items_results": result.get("items_results", []),
        "submission_preview": (submission_text or "")[:1200],
        "history_type": _normalize_history_type(history_type),
        "batch_id": str(batch_id or "").strip(),
        "batch_name": str(batch_name or "").strip(),
        "batch_size": int(batch_size or 0),
        "batch_rank": int(batch_rank or 0),
        "batch_created_at": str(batch_created_at or "").strip(),
        "score_percent": _record_percent(
            {
                "overall_score": result.get("overall_score", 0),
                "overall_out_of": result.get("overall_out_of", 0),
            }
        ),
        "manual_reviewed": False,
    })


# =========================================================
# Exports
# =========================================================
def export_single_report(record: Dict[str, Any]) -> str:
    lines = []
    lines.append("ORGANIZED GRADING REPORT")
    lines.append("=" * 90)
    lines.append(f"Title: {record.get('title', '(untitled)')}")
    lines.append(f"Generated: {record.get('timestamp', '')}")
    lines.append(f"Record ID: {record.get('id', '')}")
    lines.append(f"Overall Score: {record.get('overall_score', 0)} / {record.get('overall_out_of', 0)}")
    if record.get("batch_name"):
        lines.append(f"Batch: {record.get('batch_name')}")
    if record.get("batch_rank") and record.get("batch_size"):
        lines.append(f"Batch Rank: {record.get('batch_rank')} / {record.get('batch_size')}")
    lines.append("")
    lines.append("ITEM-BY-ITEM RESULTS")
    lines.append("-" * 90)

    for idx, item in enumerate(record.get("items_results", []), start=1):
        lines.append(f"{idx}. {item.get('name', 'Item')}")
        lines.append(f"   Origin: {item.get('item_origin', '')}")
        if item.get("description"):
            lines.append(f"   Description: {item.get('description')}")
        if item.get("mode"):
            lines.append(f"   Mode: {item.get('mode', '')}")
        if item.get("grounding"):
            lines.append(f"   Grounding: {item.get('grounding', '')}")
        lines.append(f"   Points: {item.get('earned_points', 0)} / {item.get('points', 0)}")
        if item.get("expected_answer"):
            lines.append(f"   Expected Answer / Guide: {item.get('expected_answer')}")
        if item.get("rationale"):
            lines.append(f"   Rationale: {item.get('rationale')}")
        for s in item.get("suggestions", []) or []:
            lines.append(f"   Suggestion: {s}")
        for e in item.get("evidence", [])[:3]:
            quote = (e.get("quote") or "").strip()
            source = (e.get("source") or "").strip()
            if quote:
                lines.append(f"   Evidence: {quote}" + (f" [{source}]" if source else ""))
        lines.append("")
    return "\n".join(lines)


def export_batch_report(records: List[Dict[str, Any]]) -> str:
    lines = []
    batch_name = str(records[0].get("batch_name", "")).strip() if records else ""
    lines.append("BATCH GRADING REPORT")
    lines.append("=" * 90)
    if batch_name:
        lines.append(f"Batch Name: {batch_name}")
    lines.append(f"Generated: {now_iso()}")
    lines.append(f"Number of submissions: {len(records)}")
    lines.append("")
    lines.append("RANKING")
    lines.append("-" * 90)
    for i, r in enumerate(records, start=1):
        lines.append(f"{i}. {r.get('title', '(untitled)')} -> {r.get('overall_score', 0)} / {r.get('overall_out_of', 0)}")
    lines.append("")
    lines.append("DETAILED REPORTS")
    lines.append("-" * 90)
    lines.append("")
    for r in records:
        lines.append(export_single_report(r))
        lines.append("")
        lines.append("=" * 90)
        lines.append("")
    return "\n".join(lines)


def build_docx_report(record: Dict[str, Any]) -> bytes:
    doc = Document()

    title = doc.add_heading("Grading Report", level=0)
    title.alignment = 1

    p = doc.add_paragraph()
    p.add_run("Title: ").bold = True
    p.add_run(str(record.get("title", "(untitled)")))

    p = doc.add_paragraph()
    p.add_run("Generated: ").bold = True
    p.add_run(str(record.get("timestamp", "")))

    p = doc.add_paragraph()
    p.add_run("Record ID: ").bold = True
    p.add_run(str(record.get("id", "")))

    p = doc.add_paragraph()
    p.add_run("Overall Score: ").bold = True
    p.add_run(f"{record.get('overall_score', 0)} / {record.get('overall_out_of', 0)}")

    doc.add_paragraph("")
    doc.add_heading("Item-by-Item Results", level=1)

    for idx, item in enumerate(record.get("items_results", []), start=1):
        doc.add_heading(f"{idx}. {item.get('name', 'Item')}", level=2)

        p = doc.add_paragraph()
        p.add_run("Origin: ").bold = True
        p.add_run(str(item.get("item_origin", "")))

        if item.get("description"):
            p = doc.add_paragraph()
            p.add_run("Description: ").bold = True
            p.add_run(str(item.get("description", "")))

        if item.get("mode"):
            p = doc.add_paragraph()
            p.add_run("Mode: ").bold = True
            p.add_run(str(item.get("mode", "")))

        if item.get("grounding"):
            p = doc.add_paragraph()
            p.add_run("Grounding: ").bold = True
            p.add_run(str(item.get("grounding", "")))

        p = doc.add_paragraph()
        p.add_run("Points: ").bold = True
        p.add_run(f"{item.get('earned_points', 0)} / {item.get('points', 0)}")

        if item.get("expected_answer"):
            p = doc.add_paragraph()
            p.add_run("Expected Answer / Guide: ").bold = True
            p.add_run(str(item.get("expected_answer", "")))

        if item.get("rationale"):
            p = doc.add_paragraph()
            p.add_run("Rationale: ").bold = True
            p.add_run(str(item.get("rationale", "")))

        if item.get("suggestions"):
            doc.add_paragraph("Suggestions:")
            for s in item.get("suggestions", []):
                doc.add_paragraph(str(s), style="List Bullet")

        if item.get("evidence"):
            doc.add_paragraph("Evidence:")
            for e in item.get("evidence", [])[:3]:
                quote = (e.get("quote") or "").strip()
                source = (e.get("source") or "").strip()
                if quote:
                    line = quote + (f" [{source}]" if source else "")
                    doc.add_paragraph(line, style="List Bullet")

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()


def build_batch_docx_report(records: List[Dict[str, Any]]) -> bytes:
    doc = Document()
    batch_name = str(records[0].get("batch_name", "")).strip() if records else ""

    title = doc.add_heading("Batch Grading Report", level=0)
    title.alignment = 1

    if batch_name:
        p = doc.add_paragraph()
        p.add_run("Batch Name: ").bold = True
        p.add_run(batch_name)

    p = doc.add_paragraph()
    p.add_run("Generated: ").bold = True
    p.add_run(now_iso())

    p = doc.add_paragraph()
    p.add_run("Number of submissions: ").bold = True
    p.add_run(str(len(records)))

    doc.add_heading("Ranking", level=1)
    for i, r in enumerate(records, start=1):
        doc.add_paragraph(
            f"{i}. {r.get('title', '(untitled)')} -> {r.get('overall_score', 0)} / {r.get('overall_out_of', 0)}",
            style="List Number"
        )

    doc.add_page_break()
    doc.add_heading("Detailed Reports", level=1)

    for idx, record in enumerate(records, start=1):
        doc.add_heading(f"{idx}. {record.get('title', '(untitled)')}", level=2)

        p = doc.add_paragraph()
        p.add_run("Generated: ").bold = True
        p.add_run(str(record.get("timestamp", "")))

        p = doc.add_paragraph()
        p.add_run("Overall Score: ").bold = True
        p.add_run(f"{record.get('overall_score', 0)} / {record.get('overall_out_of', 0)}")

        for j, item in enumerate(record.get("items_results", []), start=1):
            doc.add_heading(f"{j}. {item.get('name', 'Item')}", level=3)

            if item.get("description"):
                p = doc.add_paragraph()
                p.add_run("Description: ").bold = True
                p.add_run(str(item.get("description", "")))

            if item.get("mode"):
                p = doc.add_paragraph()
                p.add_run("Mode: ").bold = True
                p.add_run(str(item.get("mode", "")))

            if item.get("grounding"):
                p = doc.add_paragraph()
                p.add_run("Grounding: ").bold = True
                p.add_run(str(item.get("grounding", "")))

            p = doc.add_paragraph()
            p.add_run("Points: ").bold = True
            p.add_run(f"{item.get('earned_points', 0)} / {item.get('points', 0)}")

            if item.get("rationale"):
                p = doc.add_paragraph()
                p.add_run("Rationale: ").bold = True
                p.add_run(str(item.get("rationale", "")))

        if idx < len(records):
            doc.add_page_break()

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()


def build_single_report_html(record: Dict[str, Any]) -> str:
    rows = []
    for idx, item in enumerate(record.get("items_results", []), start=1):
        suggestions = "".join(f"<li>{s}</li>" for s in item.get("suggestions", []))
        evidence = "".join(
            f"<li>{(e.get('quote') or '').strip()}</li>"
            for e in item.get("evidence", [])[:3]
            if (e.get("quote") or "").strip()
        )

        rows.append(f"""
        <div class="item">
            <h3>{idx}. {item.get("name", "Item")}</h3>
            <div class="meta">
                <span>Origin: {item.get("item_origin", "")}</span>
                <span>Mode: {item.get("mode", "")}</span>
                <span>Grounding: {item.get("grounding", "")}</span>
                <span>Points: {item.get("earned_points", 0)} / {item.get("points", 0)}</span>
            </div>
            <p><strong>Description:</strong> {item.get("description", "")}</p>
            <p><strong>Expected Answer / Guide:</strong> {item.get("expected_answer", "")}</p>
            <p><strong>Rationale:</strong> {item.get("rationale", "")}</p>
            {"<div><strong>Suggestions:</strong><ul>" + suggestions + "</ul></div>" if suggestions else ""}
            {"<div><strong>Evidence:</strong><ul>" + evidence + "</ul></div>" if evidence else ""}
        </div>
        """)

    return f"""
    <html>
    <head>
        <meta charset="utf-8">
        <title>Grading Report</title>
        <style>
            body {{
                font-family: Arial, sans-serif;
                margin: 40px;
                color: #1c2430;
                background: #ffffff;
            }}
            .header {{
                border-bottom: 2px solid #d9e2ef;
                padding-bottom: 16px;
                margin-bottom: 24px;
            }}
            h1 {{
                margin: 0 0 10px 0;
                color: #1c2430;
            }}
            .score-box {{
                background: #f4f8fc;
                border: 1px solid #d9e2ef;
                border-radius: 12px;
                padding: 14px 16px;
                display: inline-block;
                margin-top: 8px;
                font-weight: bold;
            }}
            .item {{
                border: 1px solid #e2eaf3;
                border-radius: 14px;
                padding: 16px;
                margin-bottom: 18px;
                background: #fafcff;
            }}
            .meta {{
                display: flex;
                gap: 14px;
                flex-wrap: wrap;
                font-size: 13px;
                color: #5b6c83;
                margin-bottom: 10px;
            }}
            h3 {{
                margin-top: 0;
            }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>Grading Report</h1>
            <p><strong>Title:</strong> {record.get("title", "(untitled)")}</p>
            <p><strong>Generated:</strong> {record.get("timestamp", "")}</p>
            <p><strong>Record ID:</strong> {record.get("id", "")}</p>
            <div class="score-box">Overall Score: {record.get("overall_score", 0)} / {record.get("overall_out_of", 0)}</div>
        </div>
        {''.join(rows)}
    </body>
    </html>
    """


def build_batch_report_html(records: List[Dict[str, Any]]) -> str:
    batch_name = str(records[0].get("batch_name", "")).strip() if records else ""
    ranking = "".join(
        f"<li>{r.get('title', '(untitled)')} â€” {r.get('overall_score', 0)} / {r.get('overall_out_of', 0)}</li>"
        for r in records
    )

    blocks = []
    for r in records:
        blocks.append(f"""
        <div class="submission">
            <h2>{r.get("title", "(untitled)")}</h2>
            <p><strong>Generated:</strong> {r.get("timestamp", "")}</p>
            <div class="score-box">Score: {r.get("overall_score", 0)} / {r.get("overall_out_of", 0)}</div>
        </div>
        """)

    return f"""
    <html>
    <head>
        <meta charset="utf-8">
        <title>Batch Grading Report</title>
        <style>
            body {{
                font-family: Arial, sans-serif;
                margin: 40px;
                color: #1c2430;
                background: #ffffff;
            }}
            h1, h2 {{
                color: #1c2430;
            }}
            .section {{
                margin-bottom: 28px;
            }}
            .submission {{
                border: 1px solid #e2eaf3;
                border-radius: 14px;
                padding: 16px;
                margin-bottom: 16px;
                background: #fafcff;
            }}
            .score-box {{
                background: #f4f8fc;
                border: 1px solid #d9e2ef;
                border-radius: 12px;
                padding: 10px 14px;
                display: inline-block;
                font-weight: bold;
            }}
        </style>
    </head>
    <body>
        <div class="section">
            <h1>Batch Grading Report</h1>
            {f"<p><strong>Batch Name:</strong> {escape_html(batch_name)}</p>" if batch_name else ""}
            <p><strong>Generated:</strong> {now_iso()}</p>
            <p><strong>Number of submissions:</strong> {len(records)}</p>
        </div>

        <div class="section">
            <h2>Ranking</h2>
            <ol>{ranking}</ol>
        </div>

        <div class="section">
            <h2>Submissions</h2>
            {''.join(blocks)}
        </div>
    </body>
    </html>
    """



