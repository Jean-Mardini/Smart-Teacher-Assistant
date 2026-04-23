"""PPTX export: one full-bleed image per slide from preview screenshots (data URLs)."""

from __future__ import annotations

import base64
import re
from io import BytesIO
from typing import Sequence

_DATA_URI_RE = re.compile(
    r"^data:image/(png|jpeg|jpe|jpg|webp);base64,(.+)$",
    re.IGNORECASE | re.DOTALL,
)


def _decode_data_uri(uri: str) -> bytes:
    if not uri or not isinstance(uri, str):
        raise ValueError("Each image must be a non-empty data URL string.")
    m = _DATA_URI_RE.match(uri.strip())
    if not m:
        raise ValueError(
            'Expected data URLs like data:image/png;base64,... or data:image/jpeg;base64,...'
        )
    try:
        return base64.b64decode(m.group(2), validate=False)
    except Exception as exc:
        raise ValueError("Invalid base64 image payload.") from exc


def build_pptx_full_bleed_images(images: Sequence[str]) -> bytes:
    """Each ``data:image/...;base64,...`` becomes one slide; picture fills the slide."""
    if not images:
        raise ValueError("At least one image is required.")

    raw_list = [_decode_data_uri(u) for u in images]

    from PIL import Image
    from pptx import Presentation

    def _to_rgb(im: Image.Image) -> Image.Image:
        if im.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[-1])
            return bg
        return im.convert("RGB")

    pil_images: list[Image.Image] = []
    for raw in raw_list:
        pil_images.append(_to_rgb(Image.open(BytesIO(raw))))

    w0, h0 = pil_images[0].size
    if w0 <= 0 or h0 <= 0:
        raise ValueError("Invalid image dimensions.")

    # Web / CSS px at 96 DPI → EMU (same scaling as typical html-to-image output).
    emu_per_px = 914400 / 96.0
    slide_w = int(w0 * emu_per_px)
    slide_h = int(h0 * emu_per_px)

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = prs.slide_layouts[5]

    def _strip_placeholder_shapes(slide) -> None:
        for shape in list(slide.shapes):
            el = shape._element
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)

    for im in pil_images:
        buf = BytesIO()
        im.save(buf, format="PNG")
        buf.seek(0)
        slide = prs.slides.add_slide(blank_layout)
        _strip_placeholder_shapes(slide)
        slide.shapes.add_picture(buf, 0, 0, width=slide_w, height=slide_h)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
