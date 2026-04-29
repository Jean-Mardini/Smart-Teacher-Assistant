"""Export slide decks to PPTX — DEMO LESSON layout.

Matches the reference lesson structure:

- **Title slide**: ``DEMO LESSON`` band + deck title (+ optional subtitle).
- **Section A (e.g. overview)**: numbered headline, subtitle, **three equal rounded cards**
  in one row (``Topic — explanation`` per card), then a Notes strip.
- **Section B (e.g. location / visual)**: numbered headline, subtitle, **two columns** —
  image or gray ``Image / visual (left)`` placeholder, rounded card with bullets
  (``Topic — explanation``), then Notes.

Default content slide matches the **split** demo: **visual left + bullets right + Notes**.
``layout: grid_triple`` (from generated JSON) is the **only** case that uses three equal cards in a row.
Themes follow the user’s selected template via ``template_used`` (or ``template``) — colors, fonts, and cards.
Images use contain-fit inside the left panel so they never overlap the text column.
"""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from pptx.util import Inches, Pt


@dataclass(frozen=True)
class SlideExportTheme:
    """Visual styling aligned with ``SlideTemplate`` / Slides page writing tones."""

    page: tuple[int, int, int]
    card: tuple[int, int, int]
    line: tuple[int, int, int]
    title: tuple[int, int, int]
    sub: tuple[int, int, int]
    body: tuple[int, int, int]
    band: tuple[int, int, int]
    band_title: tuple[int, int, int]
    band_sub: tuple[int, int, int]
    callout: tuple[int, int, int]
    grad_fallback: tuple[int, int, int]
    title_font: str | None = None
    body_font: str | None = None


# Mirrors frontend ``SLIDE_TEMPLATES`` so downloaded .pptx matches the chosen tone.
SLIDE_EXPORT_THEMES: dict[str, SlideExportTheme] = {
    "academic_default": SlideExportTheme(
        page=(248, 250, 252),
        card=(255, 255, 255),
        line=(203, 213, 225),
        title=(15, 23, 42),
        sub=(71, 85, 105),
        body=(51, 65, 85),
        band=(30, 58, 95),
        band_title=(255, 255, 255),
        band_sub=(226, 232, 240),
        callout=(239, 246, 255),
        grad_fallback=(214, 234, 253),
    ),
    "minimal_clean": SlideExportTheme(
        page=(248, 250, 252),
        card=(255, 255, 255),
        line=(226, 232, 240),
        title=(30, 41, 59),
        sub=(100, 116, 139),
        body=(51, 65, 85),
        band=(71, 85, 105),
        band_title=(255, 255, 255),
        band_sub=(226, 232, 240),
        callout=(248, 250, 252),
        grad_fallback=(241, 245, 249),
    ),
    "workshop_interactive": SlideExportTheme(
        page=(240, 253, 250),
        card=(255, 255, 255),
        line=(153, 246, 228),
        title=(15, 118, 110),
        sub=(13, 148, 136),
        body=(17, 94, 89),
        band=(13, 148, 136),
        band_title=(255, 255, 255),
        band_sub=(204, 251, 241),
        callout=(236, 253, 245),
        grad_fallback=(167, 243, 208),
    ),
    "executive_summary": SlideExportTheme(
        page=(15, 23, 42),
        card=(51, 65, 85),
        line=(71, 85, 105),
        title=(254, 243, 199),
        sub=(203, 213, 225),
        body=(226, 232, 240),
        band=(245, 158, 11),
        band_title=(15, 23, 42),
        band_sub=(51, 65, 105),
        callout=(30, 41, 59),
        grad_fallback=(30, 58, 95),
        title_font="Georgia",
    ),
    "deep_technical": SlideExportTheme(
        page=(249, 250, 251),
        card=(249, 250, 251),
        line=(209, 213, 219),
        title=(17, 24, 39),
        sub=(107, 114, 128),
        body=(75, 85, 99),
        band=(17, 24, 39),
        band_title=(249, 250, 251),
        band_sub=(209, 213, 219),
        callout=(243, 244, 246),
        grad_fallback=(229, 231, 235),
        title_font="Consolas",
        body_font="Consolas",
    ),
    "story_visual": SlideExportTheme(
        page=(250, 245, 255),
        card=(255, 255, 255),
        line=(221, 214, 254),
        title=(76, 29, 149),
        sub=(124, 58, 237),
        body=(91, 33, 182),
        band=(124, 58, 237),
        band_title=(255, 255, 255),
        band_sub=(237, 233, 254),
        callout=(250, 245, 255),
        grad_fallback=(237, 233, 254),
    ),
}


def _slide_export_theme(slide_deck: dict) -> SlideExportTheme:
    # Prefer explicit ``template`` (request override), then persisted ``template_used``.
    tid = (slide_deck.get("template") or slide_deck.get("template_used") or "").strip()
    return SLIDE_EXPORT_THEMES.get(tid, SLIDE_EXPORT_THEMES["academic_default"])


def slide_deck_to_pptx_bytes(slide_deck: dict) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    def rgb(t: tuple[int, int, int]) -> RGBColor:
        return RGBColor(t[0], t[1], t[2])

    theme = _slide_export_theme(slide_deck)

    def _label_font_name() -> str:
        """Small caps band — prefer body font (sans templates), then title font."""
        if theme.body_font:
            return theme.body_font
        if theme.title_font:
            return theme.title_font
        return "Calibri"

    def _accent_band(slide, left: int, top: int, width: int) -> None:
        """Thin accent line using template ``band`` color (matches chosen tone)."""
        h = int(Pt(3))
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            int(left),
            int(top),
            int(width),
            h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(theme.band)
        try:
            bar.line.fill.background()
        except Exception:
            pass

    def _slide_bg(slide) -> None:
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = rgb(theme.page)
        except Exception:
            pass

    def _soft_shadow(shape) -> None:
        try:
            sh = shape.shadow
            sh.inherit = False
            sh.style = "outer"
            sh.distance = Pt(3)
            sh.blur_radius = Pt(6)
            sh.transparency = 0.65
        except Exception:
            pass

    def _rounded_card(slide, left, top, width, height, fill_t: tuple[int, int, int], line_t: tuple[int, int, int]):
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(left),
            int(top),
            int(width),
            int(height),
        )
        try:
            sh.adjustments[0] = 0.1
        except Exception:
            pass
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill_t)
        ln = sh.line
        ln.color.rgb = rgb(line_t)
        ln.width = Pt(1)
        _soft_shadow(sh)
        return sh

    def _image_dimensions(path: str) -> tuple[int, int] | None:
        try:
            from PIL import Image as PILImage
        except ModuleNotFoundError:
            return None
        try:
            with PILImage.open(path) as image:
                return image.size
        except Exception:
            return None

    def _add_picture_fit(slide, image_path: str, left, top, box_width, box_height) -> bool:
        if not image_path or not Path(image_path).exists():
            return False
        dimensions = _image_dimensions(image_path)
        if not dimensions:
            try:
                slide.shapes.add_picture(image_path, left, top, width=box_width)
                return True
            except Exception:
                return False
        width_px, height_px = dimensions
        if width_px <= 0 or height_px <= 0:
            return False
        image_ratio = width_px / height_px
        box_ratio = box_width / box_height
        if image_ratio >= box_ratio:
            target_width = box_width
            target_height = int(box_width / image_ratio)
        else:
            target_height = box_height
            target_width = int(box_height * image_ratio)
        left_offset = int((box_width - target_width) / 2)
        top_offset = int((box_height - target_height) / 2)
        try:
            slide.shapes.add_picture(
                image_path,
                left + left_offset,
                top + top_offset,
                width=target_width,
                height=target_height,
            )
            return True
        except Exception:
            return False

    def _add_picture_cover(slide, image_path: str, left: int, top: int, box_width: int, box_height: int) -> bool:
        """Fit image fully inside the box (object-fit: contain).

        The old cover/crop math could make ``target_width > box_width``, center with a negative
        offset, and bleed into the neighboring text column — pictures drew *above* body text.
        Contain keeps every pixel inside ``(left, top, box_width, box_height)``.
        """
        return _add_picture_fit(slide, image_path, left, top, box_width, box_height)

    def _gradient_hero(slide, left: int, top: int, width: int, height: int, seed: str) -> None:
        """Soft diagonal gradient when no bitmap is available (never a 'placeholder' text box)."""
        digest = hashlib.sha256(seed.encode("utf-8", errors="ignore")).digest()
        p, q = theme.page, theme.callout
        c1 = (
            min(255, p[0] + digest[0] % 45),
            min(255, p[1] + digest[1] % 40),
            min(255, p[2] + digest[2] % 35),
        )
        c2 = (
            min(255, q[0] + digest[3] % 45),
            min(255, q[1] + digest[4] % 40),
            min(255, q[2] + digest[5] % 35),
        )
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            int(left),
            int(top),
            int(width),
            int(height),
        )
        try:
            sh.line.fill.background()
        except Exception:
            pass
        try:
            fill = sh.fill
            fill.gradient()
            fill.gradient_angle = 115.0
            fill.gradient_stops[0].color.rgb = rgb(c1)
            fill.gradient_stops[1].color.rgb = rgb(c2)
        except Exception:
            sh.fill.solid()
            sh.fill.fore_color.rgb = rgb(theme.grad_fallback)

    def _apply_body_font(run) -> None:
        if theme.body_font:
            try:
                run.font.name = theme.body_font
            except Exception:
                pass

    def _markdown_bullets(tf, bullets: list[str], size_pt: int = 14) -> None:
        tf.clear()
        tf.word_wrap = True
        for bi, bullet in enumerate(bullets):
            text = bullet.strip()
            if not text.startswith("•"):
                text = f"• {text}"
            para = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            para.space_after = Pt(9)
            if "**" not in text:
                para.text = text
                for run in para.runs:
                    run.font.size = Pt(size_pt)
                    run.font.color.rgb = rgb(theme.body)
                    _apply_body_font(run)
                continue
            para.text = ""
            parts = re.split(r"(\*\*[^*]+\*\*)", text)
            for part in parts:
                if not part:
                    continue
                run = para.add_run()
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    run.text = part[2:-2]
                    run.font.bold = True
                else:
                    run.text = part
                    run.font.bold = False
                run.font.size = Pt(size_pt)
                run.font.color.rgb = rgb(theme.body)
                _apply_body_font(run)

    def _bullet_topic_explanation(line: str) -> tuple[str, str]:
        """Split ``Topic — explanation`` using em dash, en dash, or hyphen."""
        text = " ".join(line.split()).strip()
        for sep in (" — ", " – ", " - ", "—", "–"):
            if sep in text:
                i = text.index(sep)
                a = text[:i].strip()
                b = text[i + len(sep) :].strip()
                if a and b:
                    return a, b
        return text, ""

    def _notes_runs(tf, speaker_notes: str, notes_pt: float = 10.0) -> None:
        """``Notes:`` in italics; body in regular weight."""
        tf.clear()
        tf.word_wrap = True
        np = tf.paragraphs[0]
        np.text = ""
        np.alignment = PP_ALIGN.LEFT
        r0 = np.add_run()
        r0.text = "Notes: "
        r0.font.size = Pt(notes_pt)
        r0.font.color.rgb = rgb(theme.sub)
        r0.font.italic = True
        if theme.body_font:
            try:
                r0.font.name = theme.body_font
            except Exception:
                pass
        r1 = np.add_run()
        r1.text = speaker_notes
        r1.font.size = Pt(notes_pt)
        r1.font.color.rgb = rgb(theme.body)
        r1.font.italic = False
        _apply_body_font(r1)

    def _section_divider(slide, top: int) -> None:
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            m,
            top,
            inner_w,
            int(Pt(1.25)),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(theme.line)
        try:
            bar.line.fill.background()
        except Exception:
            pass

    def _image_path(image_refs: list[str]) -> str:
        if not image_refs:
            return ""
        meta = image_lookup.get(image_refs[0], {})
        return str(meta.get("asset_path") or meta.get("path") or "")

    def _demo_lesson_banner_and_head(slide, section_n: int, title_text: str, subtitle: str, *, top_yy: int) -> int:
        """Draw DEMO LESSON band + numbered title + subtitle; return y below subtitle block."""
        yy = top_yy
        demo_h = int(Inches(0.22))
        lbl = slide.shapes.add_textbox(m, yy, int(Inches(4.0)), demo_h)
        ltf = lbl.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lp.text = "DEMO LESSON"
        lp.alignment = PP_ALIGN.LEFT
        lp.font.size = Pt(10)
        lp.font.bold = True
        lp.font.color.rgb = rgb(theme.sub)
        try:
            lp.font.name = _label_font_name()
        except Exception:
            pass

        _accent_band(slide, m, yy + demo_h + int(Inches(0.02)), inner_w)
        yy = yy + demo_h + int(Inches(0.1))
        th = int(Inches(0.5))
        head_tb = slide.shapes.add_textbox(m, yy, inner_w, th)
        hf = head_tb.text_frame
        hf.clear()
        hp = hf.paragraphs[0]
        hp.text = f"{section_n}. {title_text}"
        hp.alignment = PP_ALIGN.LEFT
        hp.font.bold = True
        hp.font.size = Pt(22)
        hp.font.color.rgb = rgb(theme.title)
        if theme.title_font:
            try:
                hp.font.name = theme.title_font
            except Exception:
                pass

        yy = yy + th + int(Inches(0.04))
        sh = int(Inches(0.36)) if subtitle else 0
        if subtitle:
            stb = slide.shapes.add_textbox(m, yy, int(inner_w * 0.96), sh)
            sf = stb.text_frame
            sf.clear()
            sp = sf.paragraphs[0]
            sp.text = subtitle
            sp.alignment = PP_ALIGN.LEFT
            sp.font.size = Pt(12)
            sp.font.color.rgb = rgb(theme.sub)
            if theme.body_font:
                try:
                    sp.font.name = theme.body_font
                except Exception:
                    pass
            yy = yy + sh + int(Inches(0.1))
        else:
            yy = yy + int(Inches(0.08))
        return yy

    def _render_demo_lesson_three_boxes(
        slide,
        section_n: int,
        title_text: str,
        subtitle: str,
        bullets: list[str],
        speaker_notes: str,
        *,
        top_yy: int,
    ) -> None:
        """Section style 1: three equal rounded cards in one row + Notes."""
        yy = _demo_lesson_banner_and_head(slide, section_n, title_text, subtitle, top_yy=top_yy)

        notes_h = int(Inches(1.35)) if len(speaker_notes) > 420 else int(Inches(1.12))
        content_h = max(int(Inches(1.55)), slide_h - yy - notes_h - m)

        tri_gap = gap
        box_w = max(int(Inches(2.05)), (inner_w - 2 * tri_gap) // 3)
        row_w = 3 * box_w + 2 * tri_gap
        row_left = m + max(0, (inner_w - row_w) // 2)

        triple = (bullets + ["—", "—", "—"])[:3]
        pad = int(Inches(0.12))
        body_pt = 11

        for col in range(3):
            bx = row_left + col * (box_w + tri_gap)
            _rounded_card(slide, bx, yy, box_w, content_h, theme.card, theme.line)
            topic, expl = _bullet_topic_explanation(triple[col])
            topic = re.sub(r"\*\*([^*]+)\*\*", r"\1", topic).strip()
            if expl:
                expl = re.sub(r"\*\*([^*]+)\*\*", r"\1", expl).strip()
            cell_text = f"{topic} — {expl}" if expl else topic
            tb = slide.shapes.add_textbox(bx + pad, yy + pad, box_w - 2 * pad, content_h - 2 * pad)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = cell_text
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(body_pt)
                run.font.color.rgb = rgb(theme.body)
                _apply_body_font(run)

        ny = yy + content_h + int(Inches(0.08))
        if speaker_notes:
            nb = slide.shapes.add_textbox(m, ny, inner_w, notes_h)
            _notes_runs(nb.text_frame, speaker_notes, 9 if len(speaker_notes) > 500 else 10)

    def _render_demo_lesson_split(
        slide,
        section_n: int,
        title_text: str,
        subtitle: str,
        bullets: list[str],
        img_path: str,
        speaker_notes: str,
        seed: str,
        *,
        top_yy: int,
    ) -> None:
        """Section style 2: image or placeholder (left) + bullet card (right) + Notes."""
        yy = _demo_lesson_banner_and_head(slide, section_n, title_text, subtitle, top_yy=top_yy)

        notes_h = int(Inches(1.35)) if len(speaker_notes) > 420 else int(Inches(1.12))
        content_h = max(int(Inches(1.65)), slide_h - yy - notes_h - m)

        left_w = int(inner_w * 0.44)
        right_w = max(int(Inches(2.4)), inner_w - left_w - gap)
        lx, rx = m, m + left_w + gap
        pad = int(Inches(0.14))

        _rounded_card(slide, lx, yy, left_w, content_h, theme.callout, theme.line)
        _rounded_card(slide, rx, yy, right_w, content_h, theme.card, theme.line)

        pic_w = left_w - 2 * pad
        pic_h = content_h - 2 * pad

        if img_path:
            if not _add_picture_cover(slide, img_path, lx + pad, yy + pad, pic_w, pic_h):
                _gradient_hero(slide, lx + pad, yy + pad, pic_w, pic_h, seed)
        else:
            ph = slide.shapes.add_textbox(lx + pad, yy + pad, pic_w, pic_h)
            ptf = ph.text_frame
            ptf.clear()
            ptf.word_wrap = True
            pp = ptf.paragraphs[0]
            pp.text = "Image / visual (left)"
            pp.alignment = PP_ALIGN.CENTER
            pp.font.size = Pt(12)
            pp.font.color.rgb = rgb(theme.sub)
            pp.font.italic = True
            if theme.body_font:
                try:
                    pp.font.name = theme.body_font
                except Exception:
                    pass
            try:
                ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

        rtb = slide.shapes.add_textbox(rx + pad, yy + pad, right_w - 2 * pad, content_h - 2 * pad)
        _markdown_bullets(rtb.text_frame, bullets if bullets else ["—"], 12)

        ny = yy + content_h + int(Inches(0.08))
        if speaker_notes:
            nb = slide.shapes.add_textbox(m, ny, inner_w, notes_h)
            _notes_runs(nb.text_frame, speaker_notes, 9 if len(speaker_notes) > 500 else 10)

    def _use_triple_box_layout(item: dict) -> bool:
        """Only explicit ``grid_triple`` uses three cards; default matches split-screen demo."""
        return (item.get("layout") or "").strip().lower() == "grid_triple"

    def _render_demo_lesson_slide(
        slide,
        section_n: int,
        title_text: str,
        subtitle: str,
        bullets: list[str],
        img_path: str,
        speaker_notes: str,
        seed: str,
        *,
        item: dict,
    ) -> None:
        top_yy = int(Inches(0.12))
        if section_n >= 2:
            _section_divider(slide, int(Inches(0.06)))
            top_yy = int(Inches(0.2))

        if _use_triple_box_layout(item):
            _render_demo_lesson_three_boxes(
                slide,
                section_n,
                title_text,
                subtitle,
                bullets,
                speaker_notes,
                top_yy=top_yy,
            )
        else:
            _render_demo_lesson_split(
                slide,
                section_n,
                title_text,
                subtitle,
                bullets,
                img_path,
                speaker_notes,
                seed,
                top_yy=top_yy,
            )

    image_lookup = {
        item.get("image_id"): item
        for item in slide_deck.get("image_catalog", [])
        if item.get("image_id")
    }

    presentation = Presentation()
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    try:
        blank_layout = presentation.slide_layouts[6]
    except IndexError:
        blank_layout = presentation.slide_layouts[5]

    gap = int(Inches(0.18))
    m = int(Inches(0.36))
    inner_w = slide_w - 2 * m

    slides_list = slide_deck.get("slides") or []
    deck_title = slide_deck.get("title", "Presentation")
    sub0 = (slides_list[0].get("subtitle") or "").strip() if slides_list else ""

    ts = presentation.slides.add_slide(blank_layout)
    _slide_bg(ts)
    cy = int(Inches(0.2))
    lab_h = int(Inches(0.24))
    dlab = ts.shapes.add_textbox(m, cy, inner_w, lab_h)
    dtf = dlab.text_frame
    dtf.clear()
    dp = dtf.paragraphs[0]
    dp.text = "DEMO LESSON"
    dp.alignment = PP_ALIGN.LEFT
    dp.font.size = Pt(10)
    dp.font.bold = True
    dp.font.color.rgb = rgb(theme.sub)
    try:
        dp.font.name = _label_font_name()
    except Exception:
        pass
    _accent_band(ts, m, cy + lab_h + int(Inches(0.02)), inner_w)
    cy += lab_h + int(Inches(0.14))
    ttl = ts.shapes.add_textbox(m, cy, inner_w, int(Inches(0.85)))
    ttf = ttl.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = deck_title
    tp.alignment = PP_ALIGN.LEFT
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = rgb(theme.title)
    if theme.title_font:
        try:
            tp.font.name = theme.title_font
        except Exception:
            pass
    cy += int(Inches(0.95))
    if sub0:
        stb = ts.shapes.add_textbox(m, cy, int(inner_w * 0.95), int(Inches(0.45)))
        sf = stb.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = sub0
        sp.font.size = Pt(13)
        sp.font.color.rgb = rgb(theme.sub)
        if theme.body_font:
            try:
                sp.font.name = theme.body_font
            except Exception:
                pass

    for si, item in enumerate(slides_list):
        slide = presentation.slides.add_slide(blank_layout)
        _slide_bg(slide)
        title_text = item.get("slide_title", "Slide")
        subtitle = (item.get("subtitle") or "").strip()
        bullets = [b for b in (item.get("bullets") or []) if isinstance(b, str) and b.strip()][:5]
        image_refs = [
            ref
            for ref in item.get("image_refs", [])
            if ref in image_lookup and (image_lookup[ref].get("asset_path") or image_lookup[ref].get("path"))
        ][:1]
        img_path = _image_path(image_refs)
        seed = f"{title_text}|demo-lesson|{si}"
        notes_body = (item.get("speaker_notes") or "").strip()
        _render_demo_lesson_slide(
            slide,
            si + 1,
            title_text,
            subtitle,
            bullets,
            img_path,
            notes_body,
            seed,
            item=item,
        )
        if notes_body:
            slide.notes_slide.notes_text_frame.text = notes_body

    buf = BytesIO()
    presentation.save(buf)
    return buf.getvalue()
