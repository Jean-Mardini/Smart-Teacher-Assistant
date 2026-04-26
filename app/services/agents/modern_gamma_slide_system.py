"""Gamma-style slide pipeline: Groq JSON → rotating layouts → HF AI images → python-pptx / live JSON."""

from __future__ import annotations

import base64
import hashlib
import os
import re
import secrets
from concurrent.futures import ThreadPoolExecutor, as_completed
import shutil
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx.util import Inches, Pt

from app.models.modern_slide_generation import GammaSlideSpec
from app.services.agents.slide_image_generator import (
    _coerce_png_bytes,
    _generate_single_image,
    assert_live_slide_deck_payload,
    ensure_slide_image_prompt,
    is_valid_live_slide_data_url,
    live_slide_image_data_url,
    live_slide_placeholder_data_url,
    slide_image_parallel_workers,
)
from app.services.llm.groq_client import call_llm_json


def _slide_image_png_fingerprint(data_url: str) -> bytes:
    """SHA-256 of decoded PNG/file bytes — used only to spot duplicate renders (not security)."""
    try:
        if "base64," not in data_url:
            return hashlib.sha256(data_url.encode("utf-8", errors="ignore")).digest()
        raw_b64 = data_url.split("base64,", 1)[1].strip()
        raw = base64.standard_b64decode(raw_b64)
        return hashlib.sha256(raw).digest()
    except Exception:
        return hashlib.sha256(str(data_url).encode("utf-8", errors="ignore")).digest()


def _dedupe_identical_slide_images(
    *,
    deck: str,
    slides: list[dict[str, Any]],
    images_by_idx: dict[int, str],
    image_style: str,
    nslides: int,
    image_warnings: list[str],
) -> None:
    """At most a few rounds — regenerates duplicates only (no O(n³) HF loops)."""
    for _round in range(3):
        fp_to_idxs: dict[bytes, list[int]] = {}
        for i in range(nslides):
            fp = _slide_image_png_fingerprint(images_by_idx[i])
            fp_to_idxs.setdefault(fp, []).append(i)
        dup_groups = [idxs for idxs in fp_to_idxs.values() if len(idxs) >= 2]
        if not dup_groups:
            break
        for idxs in dup_groups:
            for dup_i in idxs[1:]:
                nonce = secrets.token_hex(4)
                slides[dup_i]["image_prompt"] = (
                    (slides[dup_i].get("image_prompt") or "").strip()
                    + f" Unique visualization variant {nonce}; entirely different composition from other slides in this deck."
                )
                lw_dup: list[str] = []
                if dup_i == 0:
                    images_by_idx[dup_i] = live_slide_image_data_url(
                        deck,
                        slides[dup_i],
                        image_style,
                        dup_i,
                        warnings_out=lw_dup,
                        deck_slide_count=nslides,
                    )
                else:
                    images_by_idx[dup_i] = live_slide_placeholder_data_url(dup_i)
                image_warnings.extend(lw_dup)
                image_warnings.append(
                    f"Slide {dup_i + 1}: regenerated — identical image bytes were shared with another slide."
                )


# ── Slide 0: AI hero image — image-forward layouts only (matches ``SlideRenderer`` hero branch).
HERO_LAYOUT_POOL: tuple[str, ...] = (
    "split_left",
    "split_right",
    "grid",
    "feature",
    "comparison",
    "image_dominant",
    "top_bottom",
    "highlight",
    "image_top_text_bottom",
)

# ── Slides 1..n: text-first / decorative only (no hero image shown in UI).
TEXT_LAYOUT_POOL: tuple[str, ...] = (
    "text_comparison",
    "grid_cards_2x2",
    "grid_cards_3",
    "stacked_sections",
    "feature_highlight",
    "timeline_horizontal",
    "timeline_vertical",
    "icon_list",
    "quote_highlight",
    "title_top_2_columns",
    "big_number_stats",
    "split_3_columns",
    "callout_blocks",
    "process_flow_boxes",
    "text_only_centered",
    "text_only",
)

_C_PAGE = (248, 250, 252)
_C_CARD = (255, 255, 255)
_C_LINE = (203, 213, 225)
_C_TITLE = (15, 23, 42)
_C_BODY = (51, 65, 85)
_C_BAND = (30, 58, 95)
_C_CALLOUT = (239, 246, 255)


def _bullet_count(slide: dict[str, Any]) -> int:
    return len([b for b in (slide.get("bullets") or []) if str(b).strip()])


def _content_hint_layout(slide: dict[str, Any], *, hero: bool) -> str | None:
    """Prefer a layout from content when it fits the rotation pool."""
    title = (slide.get("title") or "").lower()
    n = _bullet_count(slide)
    if hero:
        if any(k in title for k in ("process", "step", "stages", "workflow", "timeline", "how it works")):
            return "image_top_text_bottom"
        return None
    if n == 2:
        return "text_comparison"
    if n == 3:
        return "grid_cards_3"
    if n >= 4:
        return "grid_cards_2x2"
    if any(k in title for k in ("process", "step", "stages", "workflow", "cycle")):
        return "timeline_vertical" if any(
            k in title for k in ("vertical", "stack", "sequence")
        ) else "timeline_horizontal"
    if any(k in title for k in ("importance", "benefit", "impact", "why it matters", "key takeaway")):
        return "callout_blocks"
    return None


def _pick_layout_no_repeat(pool: list[str], last: str | None, hint: str | None, index: int) -> str:
    """Pick from *pool* so the choice is never *last* (previous slide layout)."""
    available = [p for p in pool if p != last]
    if not available:
        available = list(dict.fromkeys(pool))
    if hint and hint in available:
        return hint
    seed = index * 37 % len(available)
    return available[seed]


def apply_first_slide_hero_rest_text_only(slides: list[dict[str, Any]]) -> None:
    """Backward-compatible no-op — assignment is handled by :func:`assign_layouts`."""
    if not slides:
        return


def assign_layouts(slides: list[dict[str, Any]]) -> None:
    """Hero layout on slide 0; diverse text layouts on slides 1..n — **never** same layout twice in a row."""
    if not slides:
        return
    if len(slides) == 1:
        slides[0]["layout"] = slides[0]["type"] = "split_left"
        return

    hero_pool = list(HERO_LAYOUT_POOL)
    text_pool = list(TEXT_LAYOUT_POOL)
    last: str | None = None
    for i, slide in enumerate(slides):
        if i == 0:
            hint = _content_hint_layout(slide, hero=True)
            lay = _pick_layout_no_repeat(hero_pool, None, hint, i)
        else:
            hint = _content_hint_layout(slide, hero=False)
            lay = _pick_layout_no_repeat(text_pool, last, hint, i)
        slide["layout"] = slide["type"] = lay
        last = lay

    _repair_consecutive_duplicate_layouts(slides)

    for i in range(1, len(slides)):
        a = (slides[i - 1].get("layout") or "").lower()
        b = (slides[i].get("layout") or "").lower()
        if a == b:
            raise RuntimeError(
                f"Consecutive duplicate slide layout {a!r} at slides {i} and {i + 1}"
            )


def _repair_consecutive_duplicate_layouts(slides: list[dict[str, Any]]) -> None:
    """Safety net if hints/pools ever produce consecutive duplicates."""
    for i in range(1, len(slides)):
        prev_l = (slides[i - 1].get("layout") or "").lower()
        cur_l = (slides[i].get("layout") or "").lower()
        if prev_l != cur_l:
            continue
        pool = list(TEXT_LAYOUT_POOL if i >= 1 else HERO_LAYOUT_POOL)
        candidates = [p for p in pool if p.lower() != prev_l]
        if not candidates:
            candidates = list(pool)
        slides[i]["layout"] = slides[i]["type"] = candidates[i % len(candidates)]


def _truncate_title(title: str, max_len: int = 52) -> str:
    t = (title or "this topic").strip()
    if len(t) <= max_len:
        return t
    return t[: max_len - 1] + "…"


def _teaching_fillers(title: str) -> list[str]:
    tt = _truncate_title(title)
    return [
        f"**Worked example** — Walk through one situation from the reading where “{tt}” changes outcomes, naming at least two variables or actors the text emphasizes.",
        f"**Misconception check** — State one common misunderstanding about this idea, then correct it strictly using wording or logic from the source.",
        f"**Boundary condition** — Identify when this claim is not supposed to apply (limitations, caveats, or exceptions the author signals).",
        f"**Transfer task** — Ask learners to connect “{tt}” to a second example in the same document or to a closely related mechanism the text compares.",
        f"**Evidence trace** — Quote or paraphrase one concrete datum (definition, mechanism, comparison, or finding) that supports this bullet’s claim.",
        f"**So what?** — Spell out why this point matters for the chapter’s argument, method, or conclusion in one complete reasoned phrase.",
    ]


def _densify_bullets(bullets: list[str], target: int, title: str) -> list[str]:
    """Reach *target* items by splitting long bullets on sentences, then short teaching pads."""
    if target <= 0:
        return []
    out: list[str] = []
    fillers = _teaching_fillers(title)
    fi = 0
    raw = [str(b).strip() for b in bullets if str(b).strip()]
    i = 0
    while i < len(raw) and len(out) < target:
        b = raw[i]
        i += 1
        if len(b) > 100 and re.search(r"[.!?]\s+[A-Za-z*•]", b):
            chunks = re.split(r"(?<=[.!?])\s+", b)
            chunks = [c.strip() for c in chunks if c.strip()]
            if len(chunks) >= 2:
                for ch in chunks:
                    if len(out) >= target:
                        break
                    out.append(ch)
                continue
        out.append(b)
    while len(out) < target:
        out.append(fillers[fi % len(fillers)])
        fi += 1
    return out[:target]


def pad_bullets_for_layout(layout: str, bullets: list[str], title: str) -> list[str]:
    """Enough bullets to fill every card/column (live preview + :func:`build_ppt`)."""
    lay = (layout or "split_left").lower()
    raw = [str(b).strip() for b in (bullets or []) if str(b).strip()]
    if not raw:
        raw = [f"**Focus** — Anchor this slide to “{_truncate_title(title)}” using the source text."]

    if lay == "grid":
        return _densify_bullets(raw, 4, title)
    if lay == "feature":
        return _densify_bullets(raw, max(4, len(raw)), title)
    if lay == "comparison":
        return _densify_bullets(raw, max(2, len(raw)), title)
    if lay == "image_dominant":
        if len(raw) >= 4:
            return raw[:4]
        return _densify_bullets(raw, 4, title)
    if lay == "highlight":
        if len(raw) >= 3:
            return raw
        return _densify_bullets(raw, 3, title)
    if lay == "text_only":
        return raw if len(raw) >= 3 else _densify_bullets(raw, max(3, len(raw)), title)
    if lay in ("text_comparison", "title_top_2_columns"):
        return _densify_bullets(raw, max(2, len(raw)), title)
    if lay == "grid_cards_2x2":
        return _densify_bullets(raw, 4, title)
    if lay == "grid_cards_3":
        return _densify_bullets(raw, 3, title)
    if lay in ("stacked_sections", "process_flow_boxes", "timeline_vertical", "icon_list"):
        return _densify_bullets(raw, max(3, min(6, len(raw) or 3)), title)
    if lay in ("timeline_horizontal", "callout_blocks", "split_3_columns"):
        return _densify_bullets(raw, max(3, min(6, len(raw) or 3)), title)
    if lay == "feature_highlight":
        return _densify_bullets(raw, max(2, min(4, len(raw) or 2)), title)
    if lay == "quote_highlight":
        return _densify_bullets(raw, max(2, min(5, len(raw) or 2)), title)
    if lay == "big_number_stats":
        return _densify_bullets(raw, max(3, min(4, len(raw) or 3)), title)
    if lay == "text_only_centered":
        return _densify_bullets(raw, max(2, min(4, len(raw) or 2)), title)
    if lay == "image_top_text_bottom":
        return raw if len(raw) >= 3 else _densify_bullets(raw, max(3, len(raw)), title)
    # Default split / hero layouts: avoid thin decks — at least four substantive lines for PPTX export.
    if len(raw) < 4:
        return _densify_bullets(raw, 4, title)
    if len(raw) > 6:
        return raw[:6]
    return raw


def _ppt_chars_per_line(width_emu: int, font_pt: float) -> int:
    """Approximate characters per wrapped line for body text (Calibri-like)."""
    width_in = max(width_emu, 1) / 914400
    denom = max(font_pt * 0.5, 4.0)
    cpl = int(width_in * 72 / denom)
    return max(14, min(cpl, 100))


def _ppt_estimated_bullets_height_emu(
    lines: list[str],
    width_emu: int,
    font_pt: int,
    *,
    space_after_pt: int = 5,
) -> int:
    """Approximate EMU height for ``bullets_tf`` output at *width_emu* with *font_pt*."""
    if width_emu <= 0:
        width_emu = int(Inches(3))
    cpl = _ppt_chars_per_line(width_emu, font_pt)
    line_h = int((font_pt * 1.33 / 72.0) * 914400)
    gap = int((space_after_pt / 72.0) * 914400)
    nonempty = [str(x).strip() for x in lines if str(x).strip()]
    if not nonempty:
        return int(Inches(0.45))
    total_wrap_lines = 0
    for t in nonempty:
        plain = re.sub(r"\*\*([^*]+)\*\*", r"\1", t)
        plain = " ".join(plain.split())
        n = len(plain)
        wrap_lines = max(1, (n + cpl - 1) // cpl)
        total_wrap_lines += wrap_lines
    body = total_wrap_lines * line_h + max(0, len(nonempty) - 1) * gap
    return int(body * 1.07) + int(Inches(0.14))


def _ppt_text_panel_height_from_bottom(
    outer_top: int,
    content_bottom: int,
    *,
    bottom_pad: int,
    outer_h_max: int,
) -> int:
    """Shrink-wrap height for the white rounded panel (matches browser preview, not full-slide card)."""
    raw = content_bottom + bottom_pad - outer_top
    return min(max(raw, int(Inches(1.15))), outer_h_max)


def _looks_like_text_only_pill(first: str) -> bool:
    """Short category line (mirrors frontend ``looksLikeCategoryPill``)."""
    s = first.strip()
    if not s or len(s) > 44:
        return False
    if "—" in s and len(s) > 32:
        return False
    words = len(s.split())
    if words > 9:
        return False
    return True


_SLIDE_LUCIDE_ICON_RULES: tuple[tuple[re.Pattern[str], str], ...] = (
    (re.compile(r"carbon\s+dioxide|\bco2\b", re.I), "Cloud"),
    (re.compile(r"water|h2o|aqueous|moisture|droplet", re.I), "Droplets"),
    (re.compile(r"oxygen|\bo2\b|atmospheric", re.I), "AirVent"),
    (re.compile(r"sunlight|solar|\bsun\b|photon|light energy", re.I), "Sun"),
    (re.compile(r"energy|atp|nadph|power", re.I), "Zap"),
    (re.compile(r"plant|sprout|seedling|vegetation", re.I), "Sprout"),
    (re.compile(r"leaf|chlorophyll|photosynthesis", re.I), "Leaf"),
    (re.compile(r"growth|increase|surge|gain", re.I), "TrendingUp"),
    (re.compile(r"process|cycle|feedback|loop|pathway", re.I), "RefreshCw"),
    (re.compile(r"\bstep\b|phase|stage", re.I), "Hash"),
    (re.compile(r"lab|microscope|experiment|observe", re.I), "Microscope"),
    (re.compile(r"idea|concept|insight|learn", re.I), "Lightbulb"),
    # Broader topics (education, AI, tech)
    (
        re.compile(
            r"artificial intelligence|machine learning|deep learning|"
            r"\bllm\b|\bgpt\b|chatgpt|genai|neural|transformer|"
            r"\bai[-– ]?(driven|powered|based|systems?)?|\bai\b(?!\w)",
            re.I,
        ),
        "Bot",
    ),
    (re.compile(r"\bfuture\b|forecast|predict|roadmap|next-gen|tomorrow", re.I), "Rocket"),
    (re.compile(r"\bexplain|explainability|interpret|transparent|transparency|xai\b", re.I), "Eye"),
    (re.compile(r"\bdata\b|analytics|dataset|statistics|numbers", re.I), "Database"),
    (re.compile(r"\bethic|fairness|bias|trust|alignment|guardrail|safe", re.I), "Shield"),
    (re.compile(r"\bhuman\b|judgment|people|social|cultural", re.I), "Users"),
    (re.compile(r"\bstudy|research|literature|paper|citation", re.I), "BookMarked"),
    (re.compile(r"\bglobal|earth|climate|planet|international\b", re.I), "Globe"),
    (re.compile(r"\bsystem|architecture|design|component|layer", re.I), "Layers"),
)

_FALLBACK_LUCIDE_ICONS: tuple[str, ...] = (
    "Sparkles",
    "BookOpen",
    "GraduationCap",
    "Layers",
    "Target",
    "Cpu",
    "BrainCircuit",
    "Globe",
    "ShieldCheck",
    "Lightbulb",
    "Rocket",
    "CircleDot",
)


def _stable_icon_pick(title: str, bullets: list[Any]) -> str:
    blob = ((title or "").strip() + "\n" + "\n".join(str(b)[:200] for b in bullets[:5])).strip() or "slide"
    h = hashlib.md5(blob.encode("utf-8")).digest()
    idx = int.from_bytes(h[:4], "big") % len(_FALLBACK_LUCIDE_ICONS)
    return _FALLBACK_LUCIDE_ICONS[idx]


def _suggested_lucide_icon_name(slide: dict[str, Any]) -> str:
    """Keyword match → Lucide name; otherwise stable pick from title + bullets (always a valid icon key)."""
    blob = f"{slide.get('title') or ''} " + " ".join(str(b) for b in (slide.get("bullets") or []))
    for rx, name in _SLIDE_LUCIDE_ICON_RULES:
        if rx.search(blob):
            return name
    return _stable_icon_pick(str(slide.get("title") or ""), list(slide.get("bullets") or []))


def _slide_should_use_icons(slide: dict[str, Any], index: int) -> bool:
    """Concept-icon bullets on body slides — not hero (index 0); skip dense timelines/quotes."""
    if index == 0:
        return False
    lay = (slide.get("layout") or "").lower()
    bullets = slide.get("bullets") or []
    n = len([b for b in bullets if str(b).strip()])
    if n == 0 or n > 6:
        return False
    if lay == "icon_list":
        return True
    if lay in ("timeline_horizontal", "timeline_vertical", "process_flow_boxes", "quote_highlight"):
        return False
    if lay in ("grid_cards_2x2", "grid_cards_3", "callout_blocks") and n <= 4:
        return True
    if lay == "split_3_columns" and n <= 6:
        return True
    if lay in TEXT_LAYOUT_POOL or lay in HERO_LAYOUT_POOL:
        return n <= 6
    return False


def _partition_text_only_bullets(
    bullets: list[str],
) -> tuple[str | None, list[str], list[str], str | None]:
    """Category pill, two columns, footer — matches ``partitionTextOnlyBullets`` in the React app."""
    rows = [str(b).strip() for b in (bullets or []) if str(b).strip()]
    footer: str | None = None
    if len(rows) >= 5:
        footer = rows.pop()
    pill: str | None = None
    if rows and _looks_like_text_only_pill(rows[0]):
        pill = rows.pop(0)
    if not rows:
        rows = ["**Focus** — Anchor this slide to the topic using the source text."]
    mid = max(1, (len(rows) + 1) // 2)
    left_b, right_b = rows[:mid], rows[mid:]
    return pill, left_b, right_b, footer


def _system_prompt(n_slides: int, *, short_bullets: bool = False) -> str:
    if short_bullets:
        # Live + exported PPTX: syllabus-style density, one JSON string per bullet (may wrap in PowerPoint).
        bullet_rules = (
            "- Each bullets array: **4 to 6** items (prefer **5** when the source is rich). "
            "Format **`**Keyword** — explanation`** (em dash). "
            "Each bullet is **one line in JSON** but **teaching-detailed**: about **24–58 words** where the source allows — "
            "define the term, state **why it matters for the argument**, add **one concrete anchor** from the text "
            "(name, number, comparison, figure, study, quote, or method), and when helpful add **implication** "
            "(who is affected, what fails if we ignore it, or how it connects to the next slide’s idea). "
            "If the source lacks an example for a point, expand with caveats and careful inference **without inventing facts**. "
            "Do not prefix with Step 1 / Step 2. One bullet per idea — no duplicate ideas."
        )
    else:
        bullet_rules = (
            "- Each bullets array: **4 to 6** items. Each bullet ONE line, about **16–40 words**, format **Keyword — explanation** "
            "(em dash). Include **definition**, a **source-tied concrete detail**, and a **short “so what”** when possible. "
            "Do not prefix with Step 1 / Step 2. No paragraphs as a single bullet — still one continuous line per bullet."
        )
    return f"""You are an expert instructional designer. Return ONLY valid JSON (no markdown fences).

Schema:
{{"slides":[{{"title":"string","bullets":["string",...],"image_prompt":"string","type":"auto"}}]}}

IMAGE PROMPTS (critical — fed to Hugging Face text-to-image for **every** slide — no web image URLs):
- Every slide **must** include **image_prompt** describing the **exact** teaching visual for that slide only (analyze
  **title + bullets**). Images are **generated by the server** via HF Inference — phrase prompts like real art direction.
- **Preferred style words** (use liberally): realistic scientific illustration; highly detailed; educational diagram or
  educational biology illustration; clean background; precise visualization.
- Good examples:
  • "realistic scientific illustration of photosynthesis process, plant leaves absorbing sunlight, chloroplast activity, detailed biology diagram"
  • "realistic diagram of plant roots absorbing water, carbon dioxide entering leaves, sunlight on plant, educational biology illustration"
  • "plants releasing oxygen into air in sunlight, realistic nature scene, educational scientific illustration"
- Include **four or more concrete nouns** from this slide (or synonyms). No readable letters, captions, or UI in the imagined image.
- **Banned**: empty or vague prompts, cartoon mascots, unrelated abstract blobs, decorative backgrounds that ignore the topic.
- **Variation**: change composition (diagram vs scene vs macro) slide-to-slide — never reuse one generic visualization for every slide.

Rules:
- Exactly {n_slides} slides.
- Each title: clear headline (max ~14 words); prefer informative over vague.
{bullet_rules}
- type: always the string "auto".
"""


def _user_prompt(document_text: str, deck_title: str | None, n_slides: int) -> str:
    title = (deck_title or "Presentation").strip()
    return f"""Deck title: {title}

Source material (ground bullets and titles ONLY in facts implied here — do not invent unrelated topics):
---
{document_text[:100_000]}
---

Produce exactly {n_slides} slides as JSON. Bullets should read like speaker-ready notes for university or advanced secondary teaching — dense enough to export directly to slides."""


def call_structured_slides(
    document_text: str,
    n_slides: int,
    deck_title: str | None,
    *,
    short_bullets: bool = False,
) -> list[dict[str, Any]]:
    raw = call_llm_json(
        _system_prompt(n_slides, short_bullets=short_bullets),
        _user_prompt(document_text, deck_title, n_slides),
        temperature=0.4,
    )
    raw_slides = raw.get("slides") if isinstance(raw.get("slides"), list) else []
    specs: list[GammaSlideSpec] = []
    for item in raw_slides:
        if not isinstance(item, dict):
            continue
        try:
            specs.append(GammaSlideSpec.model_validate(item))
        except Exception:
            continue
    slides: list[dict[str, Any]] = []
    for spec in specs[:n_slides]:
        slides.append(
            {
                "title": spec.title.strip(),
                "bullets": spec.bullets,
                "image_prompt": spec.image_prompt.strip(),
                "type": spec.type or "auto",
            }
        )
    while len(slides) < n_slides:
        slides.append(
            {
                "title": f"Slide {len(slides) + 1}",
                "bullets": ["**Idea** — add more source text"],
                "image_prompt": "One clear educational scene that visualizes a key idea from the source text, no readable text in the image",
                "type": "auto",
            }
        )
    return slides[:n_slides]


def generate_images_for_slides(
    slides: list[dict[str, Any]],
    work_dir: Path,
    *,
    document_title: str = "Presentation",
    image_style: str = "vector_science",
) -> None:
    """Generate one PNG per slide (parallel when SLIDE_IMAGE_PARALLEL > 1)."""
    work_dir.mkdir(parents=True, exist_ok=True)
    dt = (document_title or "").strip() or "Presentation"
    style = (image_style or "").strip() or "vector_science"
    n = len(slides)
    workers = min(slide_image_parallel_workers(), n)

    def gen_one(idx: int) -> tuple[int, bytes]:
        ensure_slide_image_prompt(slides[idx], dt)
        raw, _src = _generate_single_image(slides[idx], dt, style, idx, deck_slide_count=n)
        if not raw:
            raise RuntimeError(
                f"Slide {idx + 1}: no image bytes after HF + fallbacks. Set HF_TOKEN; add PEXELS_API_KEY or OPENAI_API_KEY "
                "for SLIDE_IMAGE_FALLBACK_AFTER_HF; check HF_IMAGE_MODEL and HF quota (402)."
            )
        return idx, _coerce_png_bytes(raw)

    raw_by_idx: dict[int, bytes] = {}
    if workers <= 1 or n <= 1:
        for i in range(n):
            idx, png = gen_one(i)
            raw_by_idx[idx] = png
    else:
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = [pool.submit(gen_one, i) for i in range(n)]
            for fut in as_completed(futures):
                idx, png = fut.result()
                raw_by_idx[idx] = png

    for i in range(n):
        png = raw_by_idx[i]
        out = work_dir / f"slide_{i + 1}.png"
        out.write_bytes(png)
        slides[i]["image_path"] = str(out)
        slides[i]["image"] = str(out)


def build_ppt(slides: list[dict[str, Any]]) -> bytes:
    """Assemble PPTX bytes from slides with keys title, bullets, layout, image_path."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    def rgb(t: tuple[int, int, int]) -> RGBColor:
        return RGBColor(t[0], t[1], t[2])

    def soft_shadow(shape) -> None:
        try:
            sh = shape.shadow
            sh.inherit = False
            sh.style = "outer"
            sh.distance = Pt(3)
            sh.blur_radius = Pt(6)
            sh.transparency = 0.62
        except Exception:
            pass

    def rounded_card(slide, left, top, w, h, fill_t: tuple[int, int, int], line_t: tuple[int, int, int]):
        sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(left), int(top), int(w), int(h))
        try:
            sh.adjustments[0] = 0.1
        except Exception:
            pass
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill_t)
        ln = sh.line
        ln.color.rgb = rgb(line_t)
        ln.width = Pt(1)
        soft_shadow(sh)
        return sh

    def clear_shapes(slide) -> None:
        for shape in list(slide.shapes):
            el = shape._element
            p = el.getparent()
            if p is not None:
                p.remove(el)

    def img_dims(path: str) -> tuple[int, int] | None:
        try:
            from PIL import Image as PILImage

            with PILImage.open(path) as im:
                return im.size
        except Exception:
            return None

    def picture_cover(slide, path: str, left: int, top: int, bw: int, bh: int) -> None:
        """Fit image inside the bw×bh EMU box (like ``object-fit: contain``), centered.

        Previous "cover" scaling filled the box by enlarging past its edges; overflow was clipped by
        the slide, so photos looked zoomed/cropped. Contain keeps the full image visible with letterboxing.
        """
        if not path or not Path(path).is_file():
            return
        d = img_dims(path)
        if not d:
            slide.shapes.add_picture(path, left, top, width=bw)
            return
        wpx, hpx = d
        if wpx <= 0 or hpx <= 0:
            return
        # Scale to fit entirely inside bw × bh (same aspect ratio as pixels).
        th_fit_w = int(bw * hpx / wpx)
        if th_fit_w <= bh:
            tw, th = bw, th_fit_w
        else:
            tw, th = int(bh * wpx / hpx), bh
        x0 = left + (bw - tw) // 2
        y0 = top + (bh - th) // 2
        slide.shapes.add_picture(path, x0, y0, width=tw, height=th)

    def bullets_tf(tf, bullets: list[str], size: int = 12) -> None:
        tf.clear()
        tf.word_wrap = True
        for bi, b in enumerate(bullets):
            text = b.strip()
            if not text.startswith("•"):
                text = f"• {text}"
            para = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            para.space_after = Pt(5)
            if "**" not in text:
                para.text = text
                for r in para.runs:
                    r.font.size = Pt(size)
                    r.font.color.rgb = rgb(_C_BODY)
                continue
            para.text = ""
            parts = re.split(r"(\*\*[^*]+\*\*)", text)
            for part in parts:
                if not part:
                    continue
                run = para.add_run()
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    run.text = part[2:-2]
                    run.font.bold = True
                else:
                    run.text = part
                    run.font.bold = False
                run.font.size = Pt(size)
                run.font.color.rgb = rgb(_C_BODY)

    def title_block(slide, m: int, inner_w: int, title: str, h: int = int(Inches(0.62))) -> int:
        tb = slide.shapes.add_textbox(m, int(Inches(0.22)), inner_w, h)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = rgb(_C_TITLE)
        return int(Inches(0.22)) + h + int(Inches(0.08))

    prs = Presentation()
    # Default template is 4:3 (10×7.5"); generated slide art is overwhelmingly 16:9 — match canvas so layouts
    # align with images and viewers don’t letterbox oddly.
    prs.slide_width = int(Inches(13.333333333))
    prs.slide_height = int(Inches(7.5))
    try:
        blank = prs.slide_layouts[6]
    except IndexError:
        blank = prs.slide_layouts[5]

    sw, sh = int(prs.slide_width), int(prs.slide_height)
    m = int(Inches(0.36))
    gap = int(Inches(0.16))
    inner = sw - 2 * m

    for spec in slides:
        slide = prs.slides.add_slide(blank)
        clear_shapes(slide)
        try:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(_C_PAGE)
        except Exception:
            pass

        title = spec.get("title", "Slide")
        layout = spec.get("layout") or "split_left"
        bullets: list[str] = pad_bullets_for_layout(
            str(layout), list(spec.get("bullets") or []), str(title)
        )
        path = str(spec.get("image_path") or "")

        if layout == "image_dominant":
            y0 = int(Inches(0.16))
            body_h = sh - y0 - m
        elif layout in TEXT_LAYOUT_POOL:
            y0 = int(Inches(0.22))
            body_h = sh - y0 - m
        else:
            y0 = title_block(slide, m, inner, title)
            body_h = sh - y0 - m

        if layout == "split_left":
            tw = int(inner * 0.4)
            vw = inner - tw - gap
            tb_w = tw - int(Inches(0.32))
            est_txt = _ppt_estimated_bullets_height_emu(bullets, tb_w, 12)
            card_h = min(max(est_txt + int(Inches(0.32)), int(Inches(1.05))), body_h)
            rounded_card(slide, m, y0, tw, card_h, _C_CARD, _C_LINE)
            tb = slide.shapes.add_textbox(
                m + int(Inches(0.16)),
                y0 + int(Inches(0.14)),
                tb_w,
                card_h - int(Inches(0.28)),
            )
            bullets_tf(tb.text_frame, bullets, 12)
            picture_cover(slide, path, m + tw + gap, y0, vw, card_h)
        elif layout == "split_right":
            vw = int(inner * 0.44)
            tw = inner - vw - gap
            tb_w = tw - int(Inches(0.32))
            est_txt = _ppt_estimated_bullets_height_emu(bullets, tb_w, 12)
            card_h = min(max(est_txt + int(Inches(0.32)), int(Inches(1.05))), body_h)
            picture_cover(slide, path, m, y0, vw, card_h)
            rounded_card(slide, m + vw + gap, y0, tw, card_h, _C_CARD, _C_LINE)
            tb = slide.shapes.add_textbox(
                m + vw + gap + int(Inches(0.16)),
                y0 + int(Inches(0.14)),
                tb_w,
                card_h - int(Inches(0.28)),
            )
            bullets_tf(tb.text_frame, bullets, 12)
        elif layout == "grid":
            hero_h = int(body_h * 0.42)
            picture_cover(slide, path, m, y0, inner, hero_h)
            y1 = y0 + hero_h + gap
            h2 = sh - y1 - m
            cols, rows = 2, 2
            cw = (inner - gap) // cols
            rh = (h2 - gap) // rows
            cells = [(0, 0), (1, 0), (0, 1), (1, 1)]
            grid_b = bullets[:4]
            for ci, (cx, cy) in enumerate(cells):
                x = m + cx * (cw + gap)
                y = y1 + cy * (rh + gap)
                rounded_card(slide, x, y, cw, rh, _C_CALLOUT, _C_LINE)
                tb = slide.shapes.add_textbox(x + int(Inches(0.1)), y + int(Inches(0.08)), cw - int(Inches(0.2)), rh - int(Inches(0.16)))
                bullets_tf(tb.text_frame, [grid_b[ci]], 11)
        elif layout == "feature":
            fw = int(inner * 0.46)
            sw_ = inner - fw - gap
            rounded_card(slide, m, y0, fw, int(body_h * 0.58), _C_CALLOUT, _C_LINE)
            tb = slide.shapes.add_textbox(m + int(Inches(0.18)), y0 + int(Inches(0.14)), fw - int(Inches(0.36)), int(body_h * 0.58) - int(Inches(0.28)))
            bullets_tf(tb.text_frame, bullets[:3], 14)
            yb = y0 + int(body_h * 0.58) + gap
            hb = body_h - int(body_h * 0.58) - gap
            rounded_card(slide, m, yb, fw, hb, _C_CARD, _C_LINE)
            tb2 = slide.shapes.add_textbox(m + int(Inches(0.16)), yb + int(Inches(0.1)), fw - int(Inches(0.32)), hb - int(Inches(0.2)))
            bullets_tf(tb2.text_frame, bullets[3:], 12)
            picture_cover(slide, path, m + fw + gap, y0, sw_, body_h)
        elif layout == "highlight":
            rail = int(inner * 0.08)
            gw = inner - rail - gap
            rounded_card(slide, m, y0, rail, body_h, _C_BAND, _C_BAND)
            step_h = (body_h - gap * (max(len(bullets), 1) - 1)) // max(len(bullets), 1)
            for si, b in enumerate(bullets):
                yy = y0 + si * (step_h + gap)
                rounded_card(slide, m + rail + gap, yy, gw, step_h, _C_CARD, _C_LINE)
                tb = slide.shapes.add_textbox(m + rail + gap + int(Inches(0.12)), yy + int(Inches(0.08)), int(gw * 0.62), step_h - int(Inches(0.16)))
                bullets_tf(tb.text_frame, [b], 12)
            thumb_w = int(gw * 0.32)
            picture_cover(slide, path, m + rail + gap + int(gw * 0.64), y0, thumb_w, body_h)
        elif layout == "comparison":
            hero_h = int(body_h * 0.38)
            picture_cover(slide, path, m, y0, inner, hero_h)
            y2 = y0 + hero_h + gap
            h2 = body_h - hero_h - gap
            cw = (inner - gap) // 2
            mid = max(1, (len(bullets) + 1) // 2)
            left_b, right_b = bullets[:mid], bullets[mid:]
            rounded_card(slide, m, y2, cw, h2, _C_CARD, _C_LINE)
            rounded_card(slide, m + cw + gap, y2, cw, h2, _C_CALLOUT, _C_LINE)
            pad = int(Inches(0.14))
            lb = slide.shapes.add_textbox(m + pad, y2 + pad, cw - 2 * pad, h2 - 2 * pad)
            bullets_tf(lb.text_frame, left_b, 12)
            rb = slide.shapes.add_textbox(m + cw + gap + pad, y2 + pad, cw - 2 * pad, h2 - 2 * pad)
            bullets_tf(rb.text_frame, right_b, 12)
        elif layout == "image_dominant":
            band = int(Inches(0.72))
            rounded_card(slide, m, y0, inner, band, _C_BAND, _C_BAND)
            bt = slide.shapes.add_textbox(m + int(Inches(0.22)), y0 + int(Inches(0.12)), inner - int(Inches(0.44)), band - int(Inches(0.24)))
            tf = bt.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            img_top = y0 + band + gap
            img_h = sh - img_top - m
            picture_cover(slide, path, m, img_top, inner, img_h)
            ov_h = min(int(Inches(1.2)), int(img_h * 0.32))
            ov_t = img_top + img_h - ov_h - int(Inches(0.1))
            rounded_card(slide, m + int(Inches(0.12)), ov_t, int(inner * 0.55), ov_h, _C_CARD, _C_LINE)
            ob = slide.shapes.add_textbox(m + int(Inches(0.22)), ov_t + int(Inches(0.08)), int(inner * 0.52), ov_h - int(Inches(0.16)))
            bullets_tf(ob.text_frame, bullets[:4], 11)
        elif layout == "top_bottom":
            hero_h = int(body_h * 0.48)
            picture_cover(slide, path, m, y0, inner, hero_h)
            y1 = y0 + hero_h + gap
            h1 = sh - y1 - m
            rounded_card(slide, m, y1, inner, h1, _C_CARD, _C_LINE)
            tb = slide.shapes.add_textbox(m + int(Inches(0.2)), y1 + int(Inches(0.12)), inner - int(Inches(0.4)), h1 - int(Inches(0.24)))
            bullets_tf(tb.text_frame, bullets, 12)
        elif layout == "image_top_text_bottom":
            rem = sh - y0 - m
            hero_h = min(int(rem * 0.52), int(Inches(3.2)))
            picture_cover(slide, path, m, y0, inner, hero_h)
            y1 = y0 + hero_h + gap
            h1 = sh - y1 - m
            rounded_card(slide, m, y1, inner, h1, _C_CARD, _C_LINE)
            tb = slide.shapes.add_textbox(m + int(Inches(0.22)), y1 + int(Inches(0.14)), inner - int(Inches(0.44)), h1 - int(Inches(0.28)))
            bullets_tf(tb.text_frame, bullets, 12)
        elif layout in TEXT_LAYOUT_POOL:
            try:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = rgb((176, 214, 247))
            except Exception:
                pass
            outer_top = int(Inches(0.22))
            outer_h = sh - outer_top - m
            px = int(Inches(0.42))
            cw = inner - 2 * px

            def _ppt_title_row(cur_top: float) -> float:
                tb = slide.shapes.add_textbox(m + px, int(cur_top), cw, int(Inches(0.78)))
                ttf = tb.text_frame
                ttf.clear()
                tp = ttf.paragraphs[0]
                tp.text = title
                tp.alignment = PP_ALIGN.LEFT
                tp.font.size = Pt(26)
                tp.font.bold = True
                tp.font.color.rgb = rgb(_C_TITLE)
                return cur_top + float(Inches(0.82))

            if layout == "text_only":
                pill_t, left_b, right_b, footer_b = _partition_text_only_bullets(bullets)
                gap_cols = int(Inches(0.14))
                col_w = (cw - gap_cols) // 2
                pad_inner = int(Inches(0.18))
                inner_tw = max(1, col_w - 2 * pad_inner)
                hl = max(_ppt_estimated_bullets_height_emu(left_b, inner_tw, 11), int(Inches(0.52)))
                hr = max(_ppt_estimated_bullets_height_emu(right_b, inner_tw, 11), int(Inches(0.52)))
                cols_h_body = max(hl, hr) + 2 * pad_inner + int(Inches(0.1))
                cols_top_pre = outer_top + int(Inches(0.36))
                if pill_t:
                    cols_top_pre += int(Inches(0.4))
                cols_top_pre += int(Inches(0.76))
                footer_reserve = int(Inches(0.56)) if footer_b else int(Inches(0.26))
                slide_lim = outer_top + outer_h - footer_reserve - int(Inches(0.16))
                cols_h = min(max(cols_h_body, int(Inches(1.05))), max(1, slide_lim - cols_top_pre))
                fy_pre = cols_top_pre + cols_h + int(Inches(0.12))
                panel_bottom = fy_pre + (int(Inches(0.48)) if footer_b else int(Inches(0.2)))
                panel_h = _ppt_text_panel_height_from_bottom(
                    outer_top,
                    panel_bottom,
                    bottom_pad=int(Inches(0.32)),
                    outer_h_max=outer_h,
                )
                rounded_card(slide, m, outer_top, inner, panel_h, _C_CARD, _C_LINE)
                cur_y = outer_top + int(Inches(0.36))
                if pill_t:
                    pill_tb = slide.shapes.add_textbox(m + px, cur_y, cw, int(Inches(0.32)))
                    tf = pill_tb.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = pill_t.strip()
                    p.alignment = PP_ALIGN.LEFT
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = rgb(_C_BAND)
                    cur_y += int(Inches(0.4))
                tit_tb = slide.shapes.add_textbox(m + px, cur_y, cw, int(Inches(0.82)))
                ttf = tit_tb.text_frame
                ttf.clear()
                tp = ttf.paragraphs[0]
                tp.text = title
                tp.alignment = PP_ALIGN.LEFT
                tp.font.size = Pt(26)
                tp.font.bold = True
                tp.font.color.rgb = rgb(_C_TITLE)
                cur_y += int(Inches(0.76))
                cols_top = cur_y
                rounded_card(slide, m + px, cols_top, cw, cols_h, _C_CALLOUT, _C_LINE)
                lb = slide.shapes.add_textbox(
                    m + px + pad_inner,
                    cols_top + pad_inner,
                    inner_tw,
                    cols_h - 2 * pad_inner,
                )
                bullets_tf(lb.text_frame, left_b, 11)
                rb = slide.shapes.add_textbox(
                    m + px + col_w + gap_cols + pad_inner,
                    cols_top + pad_inner,
                    inner_tw,
                    cols_h - 2 * pad_inner,
                )
                bullets_tf(rb.text_frame, right_b, 11)
                if footer_b:
                    fy = cols_top + cols_h + int(Inches(0.12))
                    fb = slide.shapes.add_textbox(m + px, fy, cw, int(Inches(0.48)))
                    bullets_tf(fb.text_frame, [footer_b], 11)
                    for para in fb.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.color.rgb = rgb((116, 128, 148))
            elif layout == "text_comparison":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = float(outer_top + int(Inches(0.36)))
                cy = _ppt_title_row(cy)
                gap_cols = int(Inches(0.14))
                cols_top = int(cy)
                col_w = (cw - gap_cols) // 2
                pad_inner = int(Inches(0.16))
                lb_b = bullets[:1] if bullets else ["—"]
                rb_b = bullets[1:2] if len(bullets) > 1 else lb_b
                inner_tw = col_w - 2 * pad_inner
                max_h = outer_top + outer_h - cols_top - int(Inches(0.28))
                fs_tc = 12
                hl = min(
                    max(_ppt_estimated_bullets_height_emu(lb_b, inner_tw, fs_tc), int(Inches(0.62))),
                    max_h,
                )
                hr = min(
                    max(_ppt_estimated_bullets_height_emu(rb_b, inner_tw, fs_tc), int(Inches(0.62))),
                    max_h,
                )
                rounded_card(slide, m + px, cols_top, col_w, hl, _C_CALLOUT, _C_LINE)
                rounded_card(slide, m + px + col_w + gap_cols, cols_top, col_w, hr, _C_CALLOUT, _C_LINE)
                lb = slide.shapes.add_textbox(
                    m + px + pad_inner,
                    cols_top + pad_inner,
                    inner_tw,
                    hl - 2 * pad_inner,
                )
                bullets_tf(lb.text_frame, lb_b, fs_tc)
                rb = slide.shapes.add_textbox(
                    m + px + col_w + gap_cols + pad_inner,
                    cols_top + pad_inner,
                    inner_tw,
                    hr - 2 * pad_inner,
                )
                bullets_tf(rb.text_frame, rb_b, fs_tc)
            elif layout == "title_top_2_columns":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = float(outer_top + int(Inches(0.36)))
                cy = _ppt_title_row(cy)
                mid = max(1, (len(bullets) + 1) // 2)
                lb_b, rb_b = bullets[:mid], bullets[mid:]
                cols_top = int(cy)
                gap_cols = int(Inches(0.14))
                col_w = (cw - gap_cols) // 2
                pad_inner = int(Inches(0.16))
                inner_tw = col_w - 2 * pad_inner
                max_h = outer_top + outer_h - cols_top - int(Inches(0.28))
                fs_tt = 11
                hl = min(
                    max(_ppt_estimated_bullets_height_emu(lb_b, inner_tw, fs_tt), int(Inches(0.72))),
                    max_h,
                )
                hr = min(
                    max(_ppt_estimated_bullets_height_emu(rb_b, inner_tw, fs_tt), int(Inches(0.72))),
                    max_h,
                )
                rounded_card(slide, m + px, cols_top, col_w, hl, _C_CALLOUT, _C_LINE)
                rounded_card(slide, m + px + col_w + gap_cols, cols_top, col_w, hr, _C_CALLOUT, _C_LINE)
                lb = slide.shapes.add_textbox(
                    m + px + pad_inner,
                    cols_top + pad_inner,
                    inner_tw,
                    hl - 2 * pad_inner,
                )
                bullets_tf(lb.text_frame, lb_b, fs_tt)
                rb = slide.shapes.add_textbox(
                    m + px + col_w + gap_cols + pad_inner,
                    cols_top + pad_inner,
                    inner_tw,
                    hr - 2 * pad_inner,
                )
                bullets_tf(rb.text_frame, rb_b, fs_tt)
            elif layout == "grid_cards_2x2":
                title_y = float(outer_top + int(Inches(0.36)))
                y_cards_pre = int(title_y + float(Inches(0.82))) + int(Inches(0.12))
                cols = 2
                card_w = (cw - gap) // cols
                tb_inner = card_w - int(Inches(0.16))
                fs_gc = 10
                grid_b = bullets[:4]
                max_cell = int(Inches(3.15))
                heights: list[int] = []
                for ci in range(4):
                    est = _ppt_estimated_bullets_height_emu([grid_b[ci]], tb_inner, fs_gc)
                    h_cell = min(max(est + int(Inches(0.18)), int(Inches(0.52))), max_cell)
                    heights.append(h_cell)
                row0_h = max(heights[0], heights[1])
                row1_h = max(heights[2], heights[3])
                block_bottom = y_cards_pre + row0_h + gap + row1_h
                panel_h = _ppt_text_panel_height_from_bottom(
                    outer_top,
                    block_bottom,
                    bottom_pad=int(Inches(0.38)),
                    outer_h_max=outer_h,
                )
                rounded_card(slide, m, outer_top, inner, panel_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(title_y)
                y_cards = int(cy) + int(Inches(0.12))
                base_x = m + px
                y_row0 = y_cards
                rounded_card(slide, base_x, y_row0, card_w, heights[0], _C_CALLOUT, _C_LINE)
                bullets_tf(
                    slide.shapes.add_textbox(
                        base_x + int(Inches(0.08)),
                        y_row0 + int(Inches(0.06)),
                        tb_inner,
                        heights[0] - int(Inches(0.12)),
                    ).text_frame,
                    [grid_b[0]],
                    fs_gc,
                )
                rounded_card(slide, base_x + card_w + gap, y_row0, card_w, heights[1], _C_CALLOUT, _C_LINE)
                bullets_tf(
                    slide.shapes.add_textbox(
                        base_x + card_w + gap + int(Inches(0.08)),
                        y_row0 + int(Inches(0.06)),
                        tb_inner,
                        heights[1] - int(Inches(0.12)),
                    ).text_frame,
                    [grid_b[1]],
                    fs_gc,
                )
                y_row1 = y_row0 + row0_h + gap
                rounded_card(slide, base_x, y_row1, card_w, heights[2], _C_CALLOUT, _C_LINE)
                bullets_tf(
                    slide.shapes.add_textbox(
                        base_x + int(Inches(0.08)),
                        y_row1 + int(Inches(0.06)),
                        tb_inner,
                        heights[2] - int(Inches(0.12)),
                    ).text_frame,
                    [grid_b[2]],
                    fs_gc,
                )
                rounded_card(slide, base_x + card_w + gap, y_row1, card_w, heights[3], _C_CALLOUT, _C_LINE)
                bullets_tf(
                    slide.shapes.add_textbox(
                        base_x + card_w + gap + int(Inches(0.08)),
                        y_row1 + int(Inches(0.06)),
                        tb_inner,
                        heights[3] - int(Inches(0.12)),
                    ).text_frame,
                    [grid_b[3]],
                    fs_gc,
                )
            elif layout == "grid_cards_3":
                title_y = float(outer_top + int(Inches(0.36)))
                y_cards_pre = int(title_y + float(Inches(0.82))) + int(Inches(0.12))
                card_w = (cw - 2 * gap) // 3
                tb_inner = card_w - int(Inches(0.16))
                fs_g3 = 11
                slide_lim = outer_top + outer_h - int(Inches(0.22))
                max_cell_h = max(1, slide_lim - y_cards_pre)
                nshow = min(3, len(bullets))
                heights: list[int] = []
                for ci in range(nshow):
                    est = _ppt_estimated_bullets_height_emu([bullets[ci]], tb_inner, fs_g3)
                    heights.append(min(max(est + int(Inches(0.18)), int(Inches(0.52))), max_cell_h))
                row_h = max(heights) if heights else int(Inches(0.55))
                block_bottom = y_cards_pre + row_h
                panel_h = _ppt_text_panel_height_from_bottom(
                    outer_top,
                    block_bottom,
                    bottom_pad=int(Inches(0.38)),
                    outer_h_max=outer_h,
                )
                rounded_card(slide, m, outer_top, inner, panel_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(title_y)
                y_cards = int(cy) + int(Inches(0.12))
                lim2 = outer_top + panel_h - int(Inches(0.22))
                max_cell_draw = max(1, lim2 - y_cards)
                for ci in range(nshow):
                    est = _ppt_estimated_bullets_height_emu([bullets[ci]], tb_inner, fs_g3)
                    h_cell = min(max(est + int(Inches(0.18)), int(Inches(0.52))), max_cell_draw)
                    x = m + px + ci * (card_w + gap)
                    rounded_card(slide, x, y_cards, card_w, h_cell, _C_CALLOUT, _C_LINE)
                    tb = slide.shapes.add_textbox(
                        x + int(Inches(0.08)),
                        y_cards + int(Inches(0.08)),
                        tb_inner,
                        h_cell - int(Inches(0.16)),
                    )
                    bullets_tf(tb.text_frame, [bullets[ci]], fs_g3)
            elif layout == "stacked_sections":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(float(outer_top + int(Inches(0.36))))
                step_gap = int(Inches(0.12))
                block_h = max(
                    int((outer_top + outer_h - int(cy) - int(Inches(0.28))) / max(len(bullets), 1)) - step_gap,
                    int(Inches(0.85)),
                )
                for bi, b in enumerate(bullets):
                    yy = int(cy) + bi * (block_h + step_gap)
                    rounded_card(slide, m + px, yy, cw, block_h, _C_CALLOUT, _C_LINE)
                    tb = slide.shapes.add_textbox(
                        m + px + int(Inches(0.14)),
                        yy + int(Inches(0.1)),
                        cw - int(Inches(0.28)),
                        block_h - int(Inches(0.2)),
                    )
                    bullets_tf(tb.text_frame, [b], 11)
            elif layout == "feature_highlight":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cx = outer_top + int(Inches(0.4))
                tb0 = slide.shapes.add_textbox(m + px, cx, cw, int(Inches(1.1)))
                tf = tb0.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = bullets[0] if bullets else title
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = rgb(_C_TITLE)
                rest = bullets[1:] if len(bullets) > 1 else []
                tb1 = slide.shapes.add_textbox(m + px, cx + int(Inches(1.18)), cw, outer_h - int(Inches(1.9)))
                bullets_tf(tb1.text_frame, rest if rest else ["**Supporting detail** — Extend from the source."], 12)
            elif layout in ("timeline_horizontal", "process_flow_boxes"):
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(float(outer_top + int(Inches(0.36))))
                y_row = int(cy) + int(Inches(0.1))
                n = max(len(bullets), 1)
                step_w = (cw - (n - 1) * int(Inches(0.08))) // n
                step_w = max(step_w, int(Inches(1.05)))
                for ti, b in enumerate(bullets):
                    x = m + px + ti * (step_w + int(Inches(0.08)))
                    rounded_card(slide, x, y_row, step_w, int(Inches(1.35)), _C_CALLOUT, _C_LINE)
                    tb = slide.shapes.add_textbox(
                        x + int(Inches(0.06)),
                        y_row + int(Inches(0.08)),
                        step_w - int(Inches(0.12)),
                        int(Inches(1.15)),
                    )
                    bullets_tf(tb.text_frame, [b], 10)
            elif layout == "timeline_vertical":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(float(outer_top + int(Inches(0.36))))
                y0s = int(cy) + int(Inches(0.14))
                rail = int(Inches(0.14))
                step_gap = int(Inches(0.1))
                step_h = max(
                    (outer_top + outer_h - y0s - int(Inches(0.36))) // max(len(bullets), 1) - step_gap,
                    int(Inches(0.72)),
                )
                for ti, b in enumerate(bullets):
                    yy = y0s + ti * (step_h + step_gap)
                    rounded_card(slide, m + px + rail, yy, cw - rail, step_h, _C_CALLOUT, _C_LINE)
                    tb = slide.shapes.add_textbox(
                        m + px + rail + int(Inches(0.14)),
                        yy + int(Inches(0.08)),
                        cw - rail - int(Inches(0.28)),
                        step_h - int(Inches(0.16)),
                    )
                    bullets_tf(tb.text_frame, [b], 11)
            elif layout == "icon_list":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(float(outer_top + int(Inches(0.36))))
                tb = slide.shapes.add_textbox(m + px, int(cy) + int(Inches(0.08)), cw, outer_top + outer_h - int(cy) - int(Inches(0.36)))
                bullets_tf(tb.text_frame, bullets, 12)
            elif layout == "quote_highlight":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cx = outer_top + int(Inches(0.45))
                tbq = slide.shapes.add_textbox(m + px, cx, cw, int(Inches(1.35)))
                tf = tbq.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = bullets[0] if bullets else title
                p.font.size = Pt(15)
                p.font.italic = True
                p.font.color.rgb = rgb(_C_BAND)
                p.alignment = PP_ALIGN.CENTER
                tb2 = slide.shapes.add_textbox(m + px, cx + int(Inches(1.42)), cw, int(Inches(1.8)))
                bullets_tf(
                    tb2.text_frame,
                    bullets[1:] if len(bullets) > 1 else ["**Context** — Draw from the source material."],
                    11,
                )
            elif layout == "big_number_stats":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(float(outer_top + int(Inches(0.36))))
                y_row = int(cy) + int(Inches(0.14))
                n = min(len(bullets), 4)
                colw = (cw - (n - 1) * gap) // max(n, 1)
                colw = max(colw, int(Inches(0.95)))
                for ti in range(n):
                    x = m + px + ti * (colw + gap)
                    rounded_card(slide, x, y_row, colw, int(Inches(1.45)), _C_CALLOUT, _C_LINE)
                    num_tb = slide.shapes.add_textbox(x + int(Inches(0.06)), y_row + int(Inches(0.08)), colw - int(Inches(0.12)), int(Inches(0.55)))
                    ntf = num_tb.text_frame
                    ntf.clear()
                    np = ntf.paragraphs[0]
                    bt = str(bullets[ti])
                    np.text = (bt[:18] + "…") if len(bt) > 20 else bt
                    np.font.size = Pt(22)
                    np.font.bold = True
                    np.font.color.rgb = rgb(_C_BAND)
                    np.alignment = PP_ALIGN.CENTER
                    sub_tb = slide.shapes.add_textbox(
                        x + int(Inches(0.06)),
                        y_row + int(Inches(0.62)),
                        colw - int(Inches(0.12)),
                        int(Inches(0.72)),
                    )
                    bullets_tf(sub_tb.text_frame, [str(bullets[ti])], 9)
            elif layout == "split_3_columns":
                title_y = float(outer_top + int(Inches(0.36)))
                cols_top_pre = int(title_y + float(Inches(0.82))) + int(Inches(0.12))
                slide_bottom = outer_top + outer_h - int(Inches(0.2))
                max_col_h = max(1, slide_bottom - cols_top_pre)
                cw3 = (cw - 2 * gap) // 3
                third = max(1, (len(bullets) + 2) // 3)
                chunks = [
                    bullets[0:third],
                    bullets[third : third * 2],
                    bullets[third * 2 : third * 3],
                ]
                pad_x = int(Inches(0.08))
                pad_y = int(Inches(0.1))
                fs_3 = 10
                inner_w = max(1, cw3 - 2 * pad_x)
                h_cols: list[int] = []
                for ci in range(3):
                    chunk = chunks[ci][:4]
                    est = _ppt_estimated_bullets_height_emu(chunk, inner_w, fs_3)
                    h_raw = est + 2 * pad_y + int(Inches(0.04))
                    h_cols.append(min(max(h_raw, int(Inches(0.48))), max_col_h))
                block_bottom = cols_top_pre + max(h_cols)
                panel_h = _ppt_text_panel_height_from_bottom(
                    outer_top,
                    block_bottom,
                    bottom_pad=int(Inches(0.34)),
                    outer_h_max=outer_h,
                )
                rounded_card(slide, m, outer_top, inner, panel_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(title_y)
                cols_top = int(cy) + int(Inches(0.12))
                max_bottom = outer_top + panel_h - int(Inches(0.2))
                max_col_draw = max(1, max_bottom - cols_top)
                for ci in range(3):
                    x = m + px + ci * (cw3 + gap)
                    chunk = chunks[ci][:4]
                    est = _ppt_estimated_bullets_height_emu(chunk, inner_w, fs_3)
                    h_raw = est + 2 * pad_y + int(Inches(0.04))
                    h_col = min(max(h_raw, int(Inches(0.48))), max_col_draw)
                    rounded_card(slide, x, cols_top, cw3, h_col, _C_CALLOUT, _C_LINE)
                    tb_h = max(int(Inches(0.28)), h_col - 2 * pad_y)
                    tb = slide.shapes.add_textbox(x + pad_x, cols_top + pad_y, inner_w, tb_h)
                    bullets_tf(tb.text_frame, chunk, fs_3)
            elif layout == "callout_blocks":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cy = _ppt_title_row(float(outer_top + int(Inches(0.36))))
                y0b = int(cy) + int(Inches(0.1))
                ncb = len(bullets)
                bh = max((outer_top + outer_h - y0b - int(Inches(0.32))) // max(ncb, 1) - int(Inches(0.08)), int(Inches(0.72)))
                for bi, b in enumerate(bullets):
                    yy = y0b + bi * (bh + int(Inches(0.1)))
                    rounded_card(slide, m + px, yy, cw, bh, (253, 246, 227), _C_LINE)
                    tb = slide.shapes.add_textbox(m + px + int(Inches(0.14)), yy + int(Inches(0.08)), cw - int(Inches(0.28)), bh - int(Inches(0.16)))
                    bullets_tf(tb.text_frame, [b], 11)
            elif layout == "text_only_centered":
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                tb = slide.shapes.add_textbox(m + px, outer_top + int(Inches(1.05)), cw, outer_h - int(Inches(1.45)))
                bullets_tf(tb.text_frame, bullets, 12)
                for para in tb.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
            else:
                pill_t, left_b, right_b, footer_b = _partition_text_only_bullets(bullets)
                rounded_card(slide, m, outer_top, inner, outer_h, _C_CARD, _C_LINE)
                cur_y = outer_top + int(Inches(0.36))
                if pill_t:
                    pill_tb = slide.shapes.add_textbox(m + px, cur_y, cw, int(Inches(0.32)))
                    tf = pill_tb.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = pill_t.strip()
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = rgb(_C_BAND)
                    cur_y += int(Inches(0.4))
                tit_tb = slide.shapes.add_textbox(m + px, cur_y, cw, int(Inches(0.82)))
                ttf = tit_tb.text_frame
                ttf.clear()
                tp = ttf.paragraphs[0]
                tp.text = title
                tp.font.size = Pt(26)
                tp.font.bold = True
                tp.font.color.rgb = rgb(_C_TITLE)
                cur_y += int(Inches(0.76))
                gap_cols = int(Inches(0.14))
                footer_reserve = int(Inches(0.52)) if footer_b else int(Inches(0.28))
                cols_top = cur_y
                cols_h = outer_top + outer_h - cols_top - footer_reserve - int(Inches(0.22))
                cols_h = max(cols_h, int(Inches(1.35)))
                rounded_card(slide, m + px, cols_top, cw, cols_h, _C_CALLOUT, _C_LINE)
                col_w = (cw - gap_cols) // 2
                pad_inner = int(Inches(0.18))
                lb = slide.shapes.add_textbox(
                    m + px + pad_inner,
                    cols_top + pad_inner,
                    col_w - pad_inner,
                    cols_h - 2 * pad_inner,
                )
                bullets_tf(lb.text_frame, left_b, 11)
                rb = slide.shapes.add_textbox(
                    m + px + col_w + gap_cols + pad_inner,
                    cols_top + pad_inner,
                    col_w - pad_inner,
                    cols_h - 2 * pad_inner,
                )
                bullets_tf(rb.text_frame, right_b, 11)
                if footer_b:
                    fy = cols_top + cols_h + int(Inches(0.12))
                    fb = slide.shapes.add_textbox(m + px, fy, cw, int(Inches(0.5)))
                    bullets_tf(fb.text_frame, [footer_b], 11)
                    for para in fb.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.color.rgb = rgb((116, 128, 148))
        else:
            tw = int(inner * 0.4)
            vw = inner - tw - gap
            rounded_card(slide, m, y0, tw, body_h, _C_CARD, _C_LINE)
            tb = slide.shapes.add_textbox(m + int(Inches(0.16)), y0 + int(Inches(0.14)), tw - int(Inches(0.32)), body_h - int(Inches(0.28)))
            bullets_tf(tb.text_frame, bullets, 12)
            picture_cover(slide, path, m + tw + gap, y0, vw, body_h)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def run_live_slides_json(
    document_text: str,
    n_slides: int,
    deck_title: str | None,
    image_style: str,
) -> dict[str, Any]:
    """Groq (short bullets) → diverse text layouts + **no consecutive repeats**; **one hero AI image on slide 1 only**."""
    slides = call_structured_slides(document_text, n_slides, deck_title, short_bullets=True)
    assign_layouts(slides)
    for s in slides:
        lay = s.get("layout") or "split_left"
        s["bullets"] = pad_bullets_for_layout(str(lay), list(s.get("bullets") or []), str(s.get("title") or ""))
    deck = (deck_title or "").strip() or "Presentation"
    nslides = len(slides)
    workers = min(slide_image_parallel_workers(), nslides)

    image_warnings: list[str] = []

    def render_slide(i: int) -> tuple[int, str, list[str]]:
        local_warnings: list[str] = []
        if i == 0:
            img = live_slide_image_data_url(
                deck,
                slides[i],
                image_style,
                i,
                warnings_out=local_warnings,
                deck_slide_count=nslides,
            )
        else:
            img = live_slide_placeholder_data_url(i)
        return i, img, local_warnings

    # Parallel HF/Pillow calls can leave some slides with empty images on Windows; default is sequential for reliability.
    parallel_live = (os.getenv("SLIDE_IMAGE_PARALLEL_LIVE") or "").strip().lower() in (
        "1",
        "true",
        "yes",
        "on",
    )

    images_by_idx: dict[int, str] = {}
    if parallel_live and workers > 1 and nslides > 1:
        # executor.map waits for every slide image (Promise.all semantics); results follow slide order.
        with ThreadPoolExecutor(max_workers=workers) as pool:
            ordered = list(pool.map(render_slide, range(nslides)))
        for idx, img, lw in ordered:
            images_by_idx[idx] = img
            image_warnings.extend(lw)
    else:
        for i in range(nslides):
            idx, img, lw = render_slide(i)
            images_by_idx[idx] = img
            image_warnings.extend(lw)

    for i in range(nslides):
        if not is_valid_live_slide_data_url(images_by_idx.get(i, "")):
            lw_fix: list[str] = []
            if i == 0:
                images_by_idx[i] = live_slide_image_data_url(
                    deck,
                    slides[i],
                    image_style,
                    i,
                    warnings_out=lw_fix,
                    deck_slide_count=nslides,
                )
            else:
                images_by_idx[i] = live_slide_placeholder_data_url(i)
            image_warnings.extend(lw_fix)

    _dedupe_identical_slide_images(
        deck=deck,
        slides=slides,
        images_by_idx=images_by_idx,
        image_style=image_style,
        nslides=nslides,
        image_warnings=image_warnings,
    )

    out: list[dict[str, Any]] = []
    for i in range(nslides):
        lay = slides[i].get("layout") or "split_left"
        out.append(
            {
                "title": slides[i]["title"],
                "bullets": slides[i]["bullets"],
                "image_prompt": slides[i].get("image_prompt") or "",
                "type": lay,
                "layout": lay,
                "image": images_by_idx[i],
                "use_icons": _slide_should_use_icons(slides[i], i),
                "icon": _suggested_lucide_icon_name(slides[i]),
            }
        )
    assert_live_slide_deck_payload(out)
    payload: dict[str, Any] = {"deck_title": deck, "slides": out}
    if image_warnings:
        payload["warnings"] = image_warnings
    return payload


def build_ppt_from_live_slides(slides_payload: list[dict[str, Any]]) -> bytes:
    """Decode ``image`` data URLs from a live preview payload and build PPTX."""
    tmp = Path(tempfile.mkdtemp(prefix="gamma_export_"))
    try:
        norm: list[dict[str, Any]] = []
        for i, s in enumerate(slides_payload):
            img = s.get("image") or ""
            path = ""
            if isinstance(img, str) and "base64," in img:
                raw_b64 = img.split("base64,", 1)[1].strip()
                p = tmp / f"slide_{i + 1}.png"
                try:
                    p.write_bytes(base64.standard_b64decode(raw_b64))
                    path = str(p)
                except Exception:
                    path = ""
            norm.append(
                {
                    "title": s.get("title", "Slide"),
                    "bullets": list(s.get("bullets") or []),
                    "layout": (s.get("layout") or s.get("type") or "split_left"),
                    "image_path": path,
                }
            )
        return build_ppt(norm)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def run_gamma_slide_pipeline(document_text: str, n_slides: int, deck_title: str | None) -> tuple[bytes, list[str]]:
    """End-to-end: structured LLM → layouts → HF images → PPTX bytes. Returns (pptx_bytes, layouts_used)."""
    slides = call_structured_slides(document_text, n_slides, deck_title)
    assign_layouts(slides)
    for s in slides:
        lay = s.get("layout") or "split_left"
        s["bullets"] = pad_bullets_for_layout(str(lay), list(s.get("bullets") or []), str(s.get("title") or ""))
    tmp = Path(tempfile.mkdtemp(prefix="gamma_slides_"))
    try:
        img_style = (os.getenv("SLIDE_DEFAULT_IMAGE_STYLE") or "vector_science").strip()
        deck_name = (deck_title or "").strip() or "Presentation"
        generate_images_for_slides(slides, tmp, document_title=deck_name, image_style=img_style)
        pptx = build_ppt(slides)
        layouts = [s["layout"] for s in slides]
        return pptx, layouts
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
