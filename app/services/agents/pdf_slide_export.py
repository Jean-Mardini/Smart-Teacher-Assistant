"""Convert a slide-style PDF into a PPTX with one full-bleed image per page (pixel-faithful layouts)."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path


def pdf_to_pptx_bytes(pdf_path: str | Path, *, render_zoom: float = 2.0) -> bytes:
    """Each PDF page becomes one slide; the page is rendered as a single full-slide picture.

    Slide dimensions match the PDF page size so typography and layout from the source
    (e.g. *Photosynthesis — Earth's Green Engine*) are preserved.
    """
    import pymupdf

    from pptx import Presentation

    path = Path(pdf_path)
    if not path.is_file():
        raise FileNotFoundError(str(path))
    if path.suffix.lower() != ".pdf":
        raise ValueError("pdf_to_pptx_bytes expects a .pdf file")

    doc = pymupdf.open(path)
    if doc.page_count < 1:
        raise ValueError("PDF has no pages")

    first = doc[0]
    w_pt, h_pt = float(first.rect.width), float(first.rect.height)
    if w_pt <= 0 or h_pt <= 0:
        raise ValueError("Invalid PDF page dimensions")

    # OOXML: 914400 EMU per inch; PDF user units are points (1/72 inch).
    emu_per_point = 914400 / 72.0
    slide_w = int(w_pt * emu_per_point)
    slide_h = int(h_pt * emu_per_point)

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = prs.slide_layouts[5]

    mat = pymupdf.Matrix(render_zoom, render_zoom)

    def _strip_shapes(slide) -> None:
        for shape in list(slide.shapes):
            el = shape._element
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)

    for page_index in range(doc.page_count):
        page = doc[page_index]
        pix = page.get_pixmap(matrix=mat, alpha=False)
        png_io = BytesIO(pix.tobytes("png"))
        png_io.seek(0)

        if page_index == 0:
            if len(prs.slides) == 0:
                slide = prs.slides.add_slide(blank_layout)
            else:
                slide = prs.slides[0]
            _strip_shapes(slide)
        else:
            slide = prs.slides.add_slide(blank_layout)
            _strip_shapes(slide)

        slide.shapes.add_picture(png_io, 0, 0, width=slide_w, height=slide_h)

    doc.close()

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
